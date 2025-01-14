# pip install pandas
# pip install openpyxl
# pip install fsspec
# pip install Pyarrow
# pip install matplotlib
# pip install python-pptx
# pip install comtypes
# pip install PyPDF2
# pip install pywin32
# pip install requests
# pip install pyinstaller
# pip install selenium
# pip install msedge-selenium-tools selenium==3.141
# pip install --upgrade urllib3==1.26.16

import tkinter as tk
from tkinter import *
import os
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import comtypes.client
import locale
import time
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Pt, Inches
from dateutil.relativedelta import relativedelta
from PyPDF2 import PdfWriter, PdfReader
import threading
import pythoncom
import requests
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from msedge.selenium_tools import Edge, EdgeOptions
from selenium.webdriver.support.ui import Select, WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from time import sleep
from zipfile import ZipFile
from unidecode import unidecode
import pickle
from sklearn.metrics import classification_report, confusion_matrix, accuracy_score, roc_auc_score
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import StandardScaler
from sklearn.ensemble import RandomForestClassifier
from random import shuffle
from sklearn.model_selection import cross_val_score
from sklearn.model_selection import StratifiedKFold
locale.setlocale(locale.LC_ALL, 'pt_BR.UTF-8')


class Main:

    def __init__(self):
        self.root = tk.Tk()
        self.root.config(width=230, height=520)
        self.root.resizable(width=False, height=False)
        p = PhotoImage(file='Base//logo.png')
        self.root.iconphoto(False, p)
        self.root.title('Oráculo')
        self.root.config(bg='#ffffff')
        self.hoje = datetime.today()
        pythoncom.CoInitialize()

        # instancias definidas fora de __init__
        self.df = None
        self.parar = None
        self.primeiro_aviso = None
        self.segundo_aviso = None
        self.quant = None

        # Caixas
        self.caixa = tk.LabelFrame(self.root, text="Reporte", bd=5, width=20, height=20)
        self.caixa.config(bg='#ffffff')
        self.caixa.place(x=55, y=110)

        # Labels
        label = tk.Label(self.root, text="Gerador de informativo\npersonalizado")
        label.config(font=("Arial", 10), bg='#ffffff')
        label.place(x=40, y=55)
        user_label = tk.Label(self.root, text="Usuário")
        user_label.config(font=("Arial", 10), bg='#ffffff')
        user_label.place(x=50, y=293)
        senha_label = tk.Label(self.root, text="Senha")
        senha_label.config(font=("Arial", 10), bg='#ffffff')
        senha_label.place(x=50, y=338)
        taxa = tk.Label(self.root, text="Taxa de juros anual")
        taxa.config(font=("Arial", 10), bg='#ffffff')
        taxa.place(x=45, y=243)
        self.quantidade = tk.StringVar()
        self.quantidade.set("Quantidade: 0")
        self.label1 = tk.Label(self.caixa, textvariable=self.quantidade)
        self.label1.config(font=("Arial", 10), bg='#ffffff')
        self.label1.grid(row=0)
        self.counter = datetime(2000, 1, 1, 0, 0, 0)
        self.string = self.counter.strftime("Tempo: " + "%H:%M:%S")
        self.tempo = tk.Label(self.caixa)
        self.tempo.config(font=("Arial", 10), bg='#ffffff', text=self.string)
        self.tempo.grid(row=1, sticky="W")

        # Caixas de texto
        self.user = StringVar()
        self.login = tk.Entry(self.root, bg="#f2f2f2", textvariable=self.user)
        self.login.place(x=50, y=315)
        self.password = StringVar()
        self.senha = tk.Entry(self.root, bg="#f2f2f2", textvariable=self.password)
        self.senha.config(show="*")
        self.senha.place(x=50, y=360)
        self.taxa = StringVar()

        # SpinBox
        self.juros_spin = tk.Spinbox(self.root, from_=0, to=20, width=10, increment=0.5)
        self.juros_spin.delete(0, "end")
        self.juros_spin.insert(0, '8,5')
        self.juros_spin.place(x=50, y=265)

        # Botões
        self.root.bind("<Return>", (lambda event: self.funcao()))
        self.botao_gerar = tk.Button(self.root, text="Gerar Template", command=self.funcao)
        self.botao_gerar.config(width=16, bg='#ffffff', activebackground="#e6e6e6", activeforeground="Black")
        self.botao_gerar.place(x=50, y=400)
        self.botao_reiniciar = tk.Button(self.root, text="Reiniciar", command=self.ativar)
        self.botao_reiniciar.config(width=16, bg='#ffffff', activebackground="#e6e6e6", activeforeground="Black")
        self.botao_reiniciar.place(x=50, y=430)
        self.botao_sair = tk.Button(self.root, text="Sair", command=self.root.destroy)
        self.botao_sair.config(width=16, bg='#ffffff', activebackground="#e6e6e6", activeforeground="Black")
        self.botao_sair.place(x=50, y=460)

        self.root.mainloop()

    def time(self):
        # Faz o contador rodar
        self.string = self.counter.strftime("Tempo: " + "%H:%M:%S")
        self.tempo.config(text=self.string)
        self.parar = self.root.after(1000, self.time)
        self.counter += timedelta(seconds=1)

    def funcao(self):
        # Verificar se o programa esta rodando pela primeira vez e pede para reiniciar
        if self.quantidade.get() == "Quantidade: 0":
            self.senha.config(state='disabled')
            self.login.config(state='disabled')
            self.juros_spin.config(state='disabled')
            self.time()
            threading.Thread(target=self.contador).start()
        else:
            self.primeiro_aviso = 'Clique no botão Reiniciar'
            self.segundo_aviso = ' para gerar novamente'
            self.aviso()

    def contador(self):
        # Zera o contador e inicia a escolha do caminho
        self.quant = 0
        self.gerar()

    def stop(self):
        # Finaliza o programa
        self.root.after_cancel(self.parar)

    def ativar(self):
        # Reinicia o programa depois de rodar a primeira vez
        self.quantidade.set("Quantidade: 0")
        self.label1.config(text=self.quantidade)
        self.counter = datetime(2000, 1, 1, 0, 0, 0)
        self.string = self.counter.strftime("Tempo: " + "%H:%M:%S")
        self.tempo.config(text=self.string)
        self.senha.config(state='normal')
        self.login.config(state='normal')
        self.juros_spin.config(state='normal')

    # Seleciona o caminho que o programa deve seguir baseado na escolha do template
    def gerar(self):
        df_full = pd.read_excel('Base/tabela_modelo.xlsx')
        lista1 = []
        for i in df_full['DataNascimento']:
            if i != i:
                i = df_full['DataNascimento'].mean()
            lista1.append(i)
        df_full['DataNascimento'] = lista1
        lista2 = []
        for i in df_full['DataAdesao']:
            if i != i:
                i = df_full['DataAdesao'].mean()
            lista2.append(i)
        df_full['DataAdesao'] = lista2
        lista3 = []
        for i in df_full['DataAdmissao']:
            if i != i:
                i = df_full['DataAdmissao'].mean()
            lista3.append(i)
        df_full['DataAdmissao'] = lista3
        df_full['DataNascimento'] = pd.to_datetime(df_full['DataNascimento'], format="%d/%m/%Y")
        df_full['DataAdesao'] = pd.to_datetime(df_full['DataAdesao'], format="%d/%m/%Y")
        df_full['DataAdmissao'] = pd.to_datetime(df_full['DataAdmissao'], format="%d/%m/%Y")
        df_full['Idade'] = df_full['DataNascimento'].apply(lambda x:
                           (datetime.now() - relativedelta(years=x.year)).year)

        dias = []
        for x in df_full.index:
            novo = relativedelta(pd.to_datetime(df_full['DataAdesao'][x]), datetime.now())
            if novo.years == 0:
                dias.append(1)
            else:
                dias.append(novo.years)
        df_full['Anos_Adesao'] = dias

        df_full['Anos_Admissao'] = df_full['DataAdmissao'].apply(lambda x:
                                   (datetime.now() - relativedelta(years=x.year)).year)

        lista = []
        for i in df_full['UF']:
            i = unidecode(str(i).lower())
            if i == 'sao paulo' and i == '1':
                i = 'sp'
            elif i != 'sp':
                i = 'outros'
            lista.append(i)
        df_full['UF'] = lista
        df_full['UF'] = df_full['UF'].fillna(pd.Series(np.random.choice(
            list((df_full['UF'].value_counts() / df_full['UF'].count()).keys()),
            p=list((df_full['UF'].value_counts() / df_full['UF'].count()).values),
            size=len(df_full)), index=df_full.index))
        lista = []
        for i in df_full['DependentesIRF']:
            if i == i:
                i = int(i)
            lista.append(i)
        df_full['DependentesIRF'] = lista
        df_full['DependentesIRF'] = df_full['DependentesIRF'].fillna(pd.Series(np.random.choice(
            list((df_full['DependentesIRF'].value_counts() / df_full['DependentesIRF'].count()).keys()),
            p=list((df_full['DependentesIRF'].value_counts() / df_full['DependentesIRF'].count()).values),
            size=len(df_full)), index=df_full.index))
        lista = []
        for i in df_full['OpcaoIR']:
            i = str(i).lower()
            if i != 'regime regressivo' and i != 'regime progressivo':
                i = np.nan
            lista.append(i)
        df_full['OpcaoIR'] = lista
        df_full['OpcaoIR'] = df_full['OpcaoIR'].fillna(pd.Series(np.random.choice(
            list((df_full['OpcaoIR'].value_counts() / df_full['OpcaoIR'].count()).keys()),
            p=list((df_full['OpcaoIR'].value_counts() / df_full['OpcaoIR'].count()).values),
            size=len(df_full)), index=df_full.index))
        lista = []
        for i in df_full['PerfilInvestimento']:
            i = str(i).lower()
            if i == 'super conservador':
                i = 'super conser'
            elif (i != 'moderado' and i != 'conservador' and i != 'agressivo rf lp' and i != 'agressivo'
                  and i != 'super conser'):
                i = np.nan
            lista.append(i)
        df_full['PerfilInvestimento'] = lista
        df_full['PerfilInvestimento'] = df_full['PerfilInvestimento'].fillna(pd.Series(np.random.choice(
            list((df_full['PerfilInvestimento'].value_counts() / df_full['PerfilInvestimento'].count()).keys()),
            p=list((df_full['PerfilInvestimento'].value_counts() / df_full['PerfilInvestimento'].count()).values),
            size=len(df_full)), index=df_full.index))
        lista = []
        for i in df_full['Sexo']:
            i = str(i).lower()
            if i == 'masculino' or i == 'mas' or i == '1':
                i = 'm'
            elif i == 'feminino' or i == 'fem' or i == '2':
                i = 'f'
            elif i != 'm' and i != 'f':
                i = np.nan
            lista.append(i)
        df_full['Sexo'] = lista
        df_full['Sexo'] = df_full['Sexo'].fillna(pd.Series(np.random.choice(
            list((df_full['Sexo'].value_counts() / df_full['Sexo'].count()).keys()),
            p=list((df_full['Sexo'].value_counts() / df_full['Sexo'].count()).values),
            size=len(df_full)), index=df_full.index))
        lista = []
        for i in df_full['EstadoCivil']:
            i = unidecode(str(i).lower())
            if i == 'desquitado':
                i = 'divorciado'
            elif i == 'separado':
                i = 'divorciado'
            elif i == 'uniao estavel':
                i = 'uniaoestavel'
            elif i == 'marital':
                i = 'uniaoestavel'
            elif i == 'companheiro':
                i = 'uniaoestavel'
            elif i == 'outros':
                i = 'uniaoestavel'
            elif i == 'amasiado':
                i = 'uniaoestavel'
            elif i == 'naoexigido':
                i = 'solteiro'
            elif i == 'nao exigido':
                i = 'solteiro'
            elif i != 'casado' and i != 'solteiro' and i != 'viuvo' and i != 'divorciado' and i != 'uniaoestavel':
                i = np.nan
            lista.append(i)
        df_full['EstadoCivil'] = lista
        df_full['EstadoCivil'] = df_full['EstadoCivil'].fillna(pd.Series(np.random.choice(
            list((df_full['EstadoCivil'].value_counts() / df_full['EstadoCivil'].count()).keys()),
            p=list((df_full['EstadoCivil'].value_counts() / df_full['EstadoCivil'].count()).values),
            size=len(df_full)), index=df_full.index))
        lista = []
        for i in df_full['Plano']:
            i = unidecode(str(i).lower())
            if i == 'previsao':
                i = 'visao telefonica'
            elif i == 'telefonica bd':
                i = 'visao telefonica'
            elif i == 'tcoprev':
                i = 'visao telefonica'
            elif i != 'mais visao' and i != 'visao multi' and i != 'visao telefonica':
                i = np.nan
            lista.append(i)
        df_full['Plano'] = lista
        df_full['Plano'] = df_full['Plano'].fillna(pd.Series(np.random.choice(
            list((df_full['Plano'].value_counts() / df_full['Plano'].count()).keys()),
            p=list((df_full['Plano'].value_counts() / df_full['Plano'].count()).values),
            size=len(df_full)), index=df_full.index))
        df_full['UF'].replace('sp', 0, inplace=True)
        df_full['UF'].replace('outros', 1, inplace=True)
        df_full['OpcaoIR'].replace('regime progressivo', 0, inplace=True)
        df_full['OpcaoIR'].replace('regime regressivo', 1, inplace=True)
        df_full['PerfilInvestimento'].replace('conservador', 0, inplace=True)
        df_full['PerfilInvestimento'].replace('moderado', 1, inplace=True)
        df_full['PerfilInvestimento'].replace('agressivo', 2, inplace=True)
        df_full['PerfilInvestimento'].replace('super conser', 3, inplace=True)
        df_full['PerfilInvestimento'].replace('agressivo rf lp', 4, inplace=True)
        df_full['Sexo'].replace('m', 0, inplace=True)
        df_full['Sexo'].replace('f', 1, inplace=True)
        df_full['Plano'].replace('visao telefonica', 0, inplace=True)
        df_full['Plano'].replace('visao multi', 1, inplace=True)
        df_full['Plano'].replace('mais visao', 5, inplace=True)
        df_full['EstadoCivil'].replace('casado', 0, inplace=True)
        df_full['EstadoCivil'].replace('solteiro', 1, inplace=True)
        df_full['EstadoCivil'].replace('divorciado', 2, inplace=True)
        df_full['EstadoCivil'].replace('viuvo', 3, inplace=True)
        df_full['EstadoCivil'].replace('uniaoestavel', 4, inplace=True)

        filename = 'model_rfc'
        model = pickle.load(open(filename, 'rb'))
        tabela = df_full.drop(['CPF', 'Nome', 'Saldo patrocinadora', 'Saldo participante',
        'ParticipanteSA'], axis=1)[['UF', 'DependentesIRF', 'OpcaoIR', 'PerfilInvestimento', 'Sexo',
        'Plano', 'EstadoCivil', 'Idade', 'Anos_Adesao', 'Anos_Admissao']]

        prev_final = model.predict(tabela)
        prev_proba_final = model.predict_proba(tabela)

        df_full['Previsão'] = prev_final
        df_full = df_full[['CPF', 'ParticipanteSA', 'Nome', 'Previsão', 'Saldo patrocinadora', 'Saldo participante',
        'Plano', 'Idade', 'Anos_Adesao', 'OpcaoIR', 'DataAdmissao']]
        df_proba = pd.DataFrame(prev_proba_final, columns=['Probabilidade Aposentadoria', 'Probabilidade Auto/BPD',
                                                           'Probabilidade Resgate'])
        df = pd.concat([df_full, df_proba], axis=1)
        df['Plano'].replace(0, 'visao telefonica', inplace=True)
        df['Plano'].replace(1, 'visao multi', inplace=True)
        df['OpcaoIR'].replace(0, 'regime progressivo', inplace=True)
        df['OpcaoIR'].replace(1, 'regime regressivo', inplace=True)
        df = df.reset_index(drop=True)
        self.editar(df)

    def editar(self, df):
        total = []
        for i in df.index:
            total.append(df['Saldo participante'][i] + df['Saldo patrocinadora'][i])
        df['Total'] = total

        dias = []
        if datetime.today().month == 1:
            anterior = f'25/12/{datetime.today().year - 1}'
        else:
            anterior = f'25/{str(datetime.today().month - 1).zfill(2)}/{datetime.today().year}'

        for i in df.index:
            novo = relativedelta(pd.to_datetime(anterior, format='%d/%m/%Y'),
                                 pd.to_datetime(df['DataAdmissao'][i], format='%d/%m/%Y'))
            dias.append(novo.years * 12 + novo.months)
        df['Dias em Meses'] = dias

        perc = []
        for i in df.index:
            if df['Plano'][i] == 'visao multi':
                perc.append(60 if df['Dias em Meses'][i] * 0.25 >= 60 else df['Dias em Meses'][i] * 0.25)
            else:
                perc.append(3 if df['Dias em Meses'][i] <= 12 else (
                    6 if df['Dias em Meses'][i] > 12 and df['Dias em Meses'][i] <= 24 else (
                    9 if df['Dias em Meses'][i] > 24 and df['Dias em Meses'][i] <= 36 else (
                    12 if df['Dias em Meses'][i] > 36 and df['Dias em Meses'][i] <= 48 else (
                    60 if df['Dias em Meses'][i] > 48 and df['Dias em Meses'][i] <= 60 else (
                    67.5 if df['Dias em Meses'][i] > 60 and df['Dias em Meses'][i] <= 72 else (
                    75 if df['Dias em Meses'][i] > 72 and df['Dias em Meses'][i] <= 84 else (
                    82.5 if df['Dias em Meses'][i] > 84 and df['Dias em Meses'][i] <= 96 else (
                    90 if df['Dias em Meses'][i] > 96 else df['Dias em Meses'][i])))))))))
        df['Percentual de resgate'] = perc

        val = []
        for i in df.index:
            val.append(df['Saldo patrocinadora'][i] * df['Percentual de resgate'][i] / 100)
        df['Valor resgatavel'] = val

        bruto = []
        for i in df.index:
            bruto.append(df['Saldo participante'][i] + df['Valor resgatavel'][i])
        df['Valor bruto de resgate'] = bruto

        nomes = []
        for i in df['Nome'].index:
            nomes.append(df['Nome'][i].split(' ')[0].capitalize())
        df['Primeiro nome'] = nomes

        df['Tempo Plano'] = df['DataAdmissao'].apply(lambda x:
                            (datetime.now() - relativedelta(years=x.year)).year)

        lista_idade = []
        lista_plano = []
        aposentadoria = []
        for i in df.index:
            idade = 50 - df['Idade'][i]
            plano = 5 - df['Tempo Plano'][i]
            if idade > 0:
                lista_idade.append(idade)
            else:
                lista_idade.append(0)
            if plano > 0:
                lista_plano.append(plano)
            else:
                lista_plano.append(0)
        df['Tempo elegibilidade idade'] = lista_idade
        df['Tempo elegibilidade plano'] = lista_plano

        lista_sa = list(df['ParticipanteSA'])
        payment = self.valor_parcela(lista_sa)
        df['Pagamentos'] = payment

        df = df[['CPF', 'ParticipanteSA', 'Nome', 'Plano', 'Saldo participante', 'Saldo patrocinadora',
                 'Total', 'Percentual de resgate', 'Valor resgatavel', 'Valor bruto de resgate', 'Dias em Meses',
                 'Tempo elegibilidade idade', 'Tempo elegibilidade plano','Previsão', 'Pagamentos', 'Idade']]


        df_save = df[['CPF', 'ParticipanteSA', 'Nome', 'Previsão']]
        df_save.to_excel("template-resgate//previsao.xlsx", index=False)

        self.gerador(df)
        self.primeiro_aviso = 'Comunicados gerados com'
        self.segundo_aviso = ' sucesso!'
        self.aviso()
        self.stop()

    # Função que garante a formatação original dos textos
    def manter_formatacao_original(self, paragrafo, fonte_original):
        for run in paragrafo.runs:
            run.font.size = fonte_original.size
            run.font.bold = fonte_original.bold
            run.font.italic = fonte_original.italic
            run.font.color.rgb = fonte_original.color.rgb
            run.font.name = fonte_original.name
            run.font.underline = fonte_original.underline

    # Função que faz a substituição dos textos
    def substituir_texto(self, apresentacao, texto_antigo, texto_novo):
        for slide in apresentacao.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    for paragrafo in shape.text_frame.paragraphs:
                        for run in paragrafo.runs:
                            if texto_antigo in run.text:
                                fonte_original = run.font
                                run.text = run.text.replace(texto_antigo, texto_novo)
                                novo_paragrafo = shape.text_frame.add_paragraph()
                                self.manter_formatacao_original(novo_paragrafo, fonte_original)

    # Função que busca as chaves dentro o PPT
    def principal(self, cpf, valor, parcela, com, sem, renda, patrocina, percentual, resgate, bruto, id, ppt, saida,
                  aposentado, plano, dez, quinze, vinte, taxa, tag, pdf, dez1, quinze1, vinte1, dez2, quinze2, vinte2,
                  meia, cinquenta, text, payment, v10, v101, v102):
        apresentacao = Presentation(ppt)
        ppt_out = f'{saida}{id}.pptx'

        self.substituir_texto(apresentacao, '{renda}', renda)
        self.substituir_texto(apresentacao, '{patrocinadora}', patrocina)
        self.substituir_texto(apresentacao, '{valor}', valor)
        self.substituir_texto(apresentacao, '{dez}', dez)
        self.substituir_texto(apresentacao, '{quinze}', quinze)
        self.substituir_texto(apresentacao, '{vinte}', vinte)
        self.substituir_texto(apresentacao, '{percentual}', percentual)
        self.substituir_texto(apresentacao, '{resgate}', resgate)
        self.substituir_texto(apresentacao, '{bruto}', bruto)
        self.substituir_texto(apresentacao, '{plano}', plano)
        self.substituir_texto(apresentacao, '{taxa}', taxa)
        self.substituir_texto(apresentacao, '{parcela}', parcela)
        self.substituir_texto(apresentacao, '{com}', com)
        self.substituir_texto(apresentacao, '{sem}', sem)
        self.substituir_texto(apresentacao, '{dez1}', dez1)
        self.substituir_texto(apresentacao, '{quinze1}', quinze1)
        self.substituir_texto(apresentacao, '{vinte1}', vinte1)
        self.substituir_texto(apresentacao, '{dez2}', dez2)
        self.substituir_texto(apresentacao, '{quinze2}', quinze2)
        self.substituir_texto(apresentacao, '{vinte2}', vinte2)
        self.substituir_texto(apresentacao, '{meia}', meia)
        self.substituir_texto(apresentacao, '{cinquenta}', cinquenta)
        self.substituir_texto(apresentacao, '{texto}', text)

        apresentacao.save(ppt_out)
        self.imagens(ppt_out, cpf, saida, aposentado, tag, pdf, id, payment, v10, v101, v102)

    # Função que insere as imagens no PPT
    def imagens(self, ppt_out, cpf, saida, aposentado, tag, pdf, id, payment, v10, v101, v102):
        prs = Presentation(ppt_out)
        slide = prs.slides[1]
        img_path = f"{saida}foto.png"
        left = Inches(6.5)
        top = Inches(1.5)
        width = Inches(13.5)
        height = Inches(10)
        slide.shapes.add_picture(img_path, left, top, width, height)
        os.remove(f"{saida}foto.png")

        slide = prs.slides[2]
        img_path1 = f"{saida}foto1.png"
        left1 = Inches(0.4)
        top1 = Inches(5.3)
        width1 = Inches(7.2)
        height1 = Inches(5)
        img_path2 = f"{saida}foto2.png"
        left2 = Inches(6)
        img_path3 = f"{saida}foto3.png"
        left3 = Inches(11.7)
        slide.shapes.add_picture(img_path1, left1, top1, width1, height1)
        slide.shapes.add_picture(img_path2, left2, top1, width1, height1)
        slide.shapes.add_picture(img_path3, left3, top1, width1, height1)
        slide = prs.slides[3]
        img_path1 = f"{saida}foto4.png"
        left1 = Inches(0.4)
        top1 = Inches(5.3)
        width1 = Inches(7.2)
        height1 = Inches(5)
        img_path2 = f"{saida}foto5.png"
        left2 = Inches(6)
        img_path3 = f"{saida}foto6.png"
        left3 = Inches(11.7)
        slide.shapes.add_picture(img_path1, left1, top1, width1, height1)
        slide.shapes.add_picture(img_path2, left2, top1, width1, height1)
        slide.shapes.add_picture(img_path3, left3, top1, width1, height1)
        slide = prs.slides[4]
        img_path1 = f"{saida}foto7.png"
        left1 = Inches(0.4)
        top1 = Inches(5.3)
        width1 = Inches(7.2)
        height1 = Inches(5)
        img_path2 = f"{saida}foto8.png"
        left2 = Inches(6)
        img_path3 = f"{saida}foto9.png"
        left3 = Inches(11.7)
        slide.shapes.add_picture(img_path1, left1, top1, width1, height1)
        slide.shapes.add_picture(img_path2, left2, top1, width1, height1)
        slide.shapes.add_picture(img_path3, left3, top1, width1, height1)
        if tag == 2:
            if aposentado != 0:
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                xml_slides.remove(slides[6])
            else:
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                xml_slides.remove(slides[7])
        else:
            if aposentado != 0:
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                xml_slides.remove(slides[8])
            else:
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                xml_slides.remove(slides[9])

        if v10 <= 0 and v101 <= 0 and v102 <= 0 and payment != 0:
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[2])
            xml_slides.remove(slides[3])
            xml_slides.remove(slides[4])

        if v10 <= 0 and v101 <= 0 and v102 > 0 and payment != 0:
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[2])
            xml_slides.remove(slides[3])

        if v10 <= 0 and v101 > 0 and v102 > 0 and payment != 0:
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[2])

        if v10 > 0 and payment == 0:
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[3])
            xml_slides.remove(slides[4])

        if v10 <= 0 and payment == 0:
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[2])
            xml_slides.remove(slides[3])
            xml_slides.remove(slides[4])

        os.remove(f"{saida}foto1.png")
        os.remove(f"{saida}foto2.png")
        os.remove(f"{saida}foto3.png")
        os.remove(f"{saida}foto4.png")
        os.remove(f"{saida}foto5.png")
        os.remove(f"{saida}foto6.png")
        os.remove(f"{saida}foto7.png")
        os.remove(f"{saida}foto8.png")
        os.remove(f"{saida}foto9.png")

        prs.save(f'{saida}{id}.pptx')
        time.sleep(1)
        self.ppttopdf(cpf, pdf, id)

    # Função que transforma o PPT em PDF
    def ppttopdf(self, cpf, pdf, id):
        entrada = f'{pdf}{id}.pptx'
        saida = f'{pdf}{id}.pdf'
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application", pythoncom.CoInitialize())
        deck = powerpoint.Presentations.Open(entrada)
        deck.ExportAsFixedFormat(saida, 32)
        deck.Close()
        powerpoint.Quit()
        sleep(1)
        self.encriptar(cpf, saida, entrada)

    # Função que faz a encriptação do PDF
    def encriptar(self, cpf, saida, entrada):
        out = PdfWriter()
        file = PdfReader(saida)
        num = len(file.pages)
        os.remove(entrada)
        for idx in range(num):
            page = file.pages[idx]
            out.add_page(page)
        password = cpf[-4:]
        out.encrypt(password)
        with open(saida, "wb") as f:
            out.write(f)

    # Função que coleta as variáveis da tabela e cria os gráficos das 3 diferentes opções para inserir no PPT
    def gerador(self, base):
        # Primeira opção dos resgates
        for i in base.index:
            if base['Previsão'][i] == 'Resgate':
                template = 'Base//Resgate ou portabilidade.pptx'
                pdf = f'{os.path.abspath("template-resgate")}\\'
                saida = "template-resgate//"
                tag = 1
            elif base['Previsão'][i] == 'Aposentadoria':
                template = 'Base//Aposentadoria.pptx'
                pdf = f'{os.path.abspath("templates-aposentado")}\\'
                saida = "templates-aposentado//"
                tag = 2
            else:
                template = 'Base//Resgate ou portabilidade.pptx'
                pdf = f'{os.path.abspath("templates-bpd-auto")}\\'
                saida = "templates-bpd-auto//"
                tag = 3

            # Coleta das variáveis
            cpf = str(base['CPF'][i])
            id = str(base['ParticipanteSA'][i])
            renda = locale.currency(base['Saldo participante'][i], symbol=True, grouping=True,
                                    international=False)
            patrocina = locale.currency(base['Saldo patrocinadora'][i], symbol=True, grouping=True,
                                        international=False)
            valor = locale.currency(base['Total'][i], symbol=True, grouping=True, international=False)
            percentual = str(base['Percentual de resgate'][i]) + '%'
            resgate = locale.currency(base['Valor resgatavel'][i], symbol=True, grouping=True, international=False)
            bruto = locale.currency(base['Valor bruto de resgate'][i], symbol=True, grouping=True,
                                    international=False)
            plano = base['Plano'][i]
            payment = base['Pagamentos'][i]
            parcela = locale.currency(payment, symbol=True, grouping=True, international=False)
            meia = locale.currency(payment / 2, symbol=True, grouping=True, international=False)

            if base['Tempo elegibilidade idade'][i] != 0:
                if base['Tempo elegibilidade plano'][i] > base['Tempo elegibilidade idade'][i]:
                    aposentado = 0
                    taxa = str(self.juros_spin.get())
                    tx_ano = float(taxa.replace(',', '.')) / 100
                    tx = (1 + tx_ano) ** (1 / 12) - 1
                    cap = base['Total'][i]
                    tem = base['Tempo elegibilidade plano'][i]
                    par = payment
                    if par != 0:
                        tot = par * ((1 + tx) ** tem - 1) / tx + cap * (1 + tx) ** tem
                        com = locale.currency(tot * 0.005, symbol=True, grouping=True, international=False)
                        sem = locale.currency(tot / 260, symbol=True, grouping=True, international=False)
                        cinquenta = locale.currency(tot, symbol=True, grouping=True, international=False)

                    else:
                        tot = cap * (1 + tx) ** tem
                        com = locale.currency(tot * 0.005, symbol=True, grouping=True, international=False)
                        sem = locale.currency(tot / 260, symbol=True, grouping=True, international=False)
                        cinquenta = locale.currency(tot, symbol=True, grouping=True, international=False)

                else:
                    aposentado = 0
                    taxa = str(self.juros_spin.get())
                    tx_ano = float(taxa.replace(',', '.')) / 100
                    tx = (1 + tx_ano) ** (1 / 12) - 1
                    cap = base['Total'][i]
                    tem = base['Tempo elegibilidade idade'][i]
                    par = payment
                    if par != 0:
                        tot = par * ((1 + tx) ** tem - 1) / tx + cap * (1 + tx) ** tem
                        com = locale.currency(tot * 0.005, symbol=True, grouping=True, international=False)
                        sem = locale.currency(tot / 260, symbol=True, grouping=True, international=False)
                        cinquenta = locale.currency(tot, symbol=True, grouping=True, international=False)

                    else:
                        tot = cap * (1 + tx) ** tem
                        com = locale.currency(tot * 0.005, symbol=True, grouping=True, international=False)
                        sem = locale.currency(tot / 260, symbol=True, grouping=True, international=False)
                        cinquenta = locale.currency(tot, symbol=True, grouping=True, international=False)

            else:
                aposentado = 1
                cap = base['Total'][i]
                com = locale.currency(cap * 0.005, symbol=True, grouping=True, international=False)
                sem = locale.currency(cap / 260, symbol=True, grouping=True, international=False)
                cinquenta = locale.currency(base['Total'][i], symbol=True, grouping=True, international=False)

            # Criação da primeira imagem em gráfico pizza
            fig = plt.figure()
            ax = fig.add_subplot(111)
            explode = (0.01, 0.01)
            v1 = list(base[['Saldo participante', 'Saldo patrocinadora']].iloc[i])[0]
            v2 = list(base[['Saldo participante', 'Saldo patrocinadora']].iloc[i])[1]
            ax = base[['Saldo participante', 'Saldo patrocinadora']].iloc[i].plot.pie(
                colors=["#58aa7a", "#065d70"], shadow=True, explode=explode, labels=None,
                autopct=lambda p: locale.currency(p * (v1 + v2) / 100, symbol=True, grouping=True,
                                    international=False), textprops={'color': "w", 'weight': 'bold'})
            ax.yaxis.set_visible(False)
            ax.legend(loc='lower right', fontsize='x-small', labels=['Participante', 'Patrocinadora'],
                      frameon=False)
            plt.savefig(f'{saida}foto.png', transparent=True)
            plt.clf()
            plt.close()


            # Calculo do juros compostos para os 3 gráficos do terceiro slide
            v3 = v1 + v2
            taxa = str(self.juros_spin.get())
            tx = float(taxa.replace(',','.')) / 100
            tx_mes = (1 + tx) ** (1 / 12) - 1

            if plano == 'visao multi':
                plano = 'Visão Multi'
                text = 'Por não ter fins lucrativos, a taxa administrativa é\ncompetitiva em relação ao mercado. Para conferir a tabela\nvigente, consulte o site da Entidade.'
                pri = 55.00 / 12
                seg = (1 + 0.0069) ** (1 / 12) - 1
                ter = (1 + 0.0046) ** (1 / 12) - 1
                qua = (1 + 0.0028) ** (1 / 12) - 1
                qui = (1 + 0.0014) ** (1 / 12) - 1
                sex = 1900.00 / 12
            else:
                plano = 'Visão Telefônica'
                text = 'Ao se tornar aposentar pelo plano,\nnão haverá cobrança de qualquer taxa\nadministrativa.'
                pri = 93.00 / 12
                seg = (1 + 0.0102) ** (1 / 12) - 1
                ter = (1 + 0.0092) ** (1 / 12) - 1
                qua = (1 + 0.0077) ** (1 / 12) - 1
                qui = (1 + 0.0037) ** (1 / 12) - 1
                sex = 3140.00 / 12

            value = v3
            for j in range(1, 121):
                if value <= 10000:
                    value = value + value * tx_mes - pri
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1)
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1)
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1)
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1)
                elif value > 1250000:
                    value = value + value * tx_mes - sex
            v10 = value - v3

            value = v3
            for j in range(1, 181):
                if value <= 10000:
                    value = value + value * tx_mes - pri
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1)
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1)
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1)
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1)
                elif value > 1250000:
                    value = value + value * tx_mes - sex
            v15 = value - v3

            value = v3
            for j in range(1, 241):
                if value <= 10000:
                    value = value + value * tx_mes - pri
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1)
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1)
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1)
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1)
                elif value > 1250000:
                    value = value + value * tx_mes - sex
            v20 = value - v3

            # Criação do primeiro gráfico de dez anos
            ax = pd.DataFrame([[v3, v10]],columns=['Contribuição','Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v3, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v10, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right',
                      fontsize='x-small', frameon=False)
            ax.set_ylim(top=v3 + v20)
            plt.savefig(f'{saida}foto1.png', transparent=True)
            plt.clf()
            plt.close()

            # Criação do segundo gráfico de 15 anos
            ax = pd.DataFrame([[v3, v15]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v3, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v15, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right',
                      fontsize='x-small', frameon=False)
            ax.set_ylim(top=v3 + v20)
            plt.savefig(f'{saida}foto2.png', transparent=True)
            plt.clf()
            plt.close()

            # Criação do terceiro gráfico de 15 anos
            ax = pd.DataFrame([[v3, v20]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v3, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v20, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right',
                      fontsize='x-small', frameon=False)
            ax.set_ylim(top=v3 + v20)
            plt.savefig(f'{saida}foto3.png', transparent=True)
            plt.clf()
            plt.close()

            # Criação das variáveis de 10, 15 e 20 anos, e as de valor bruto, imposto e valor líquido
            dez = locale.currency(v3 + v10, symbol=True, grouping=True, international=False)
            quinze = locale.currency(v3 + v15, symbol=True, grouping=True, international=False)
            vinte = locale.currency(v3 + v20, symbol=True, grouping=True, international=False)

            ### calculando os juros compostos com metade da contribuição
            v7 = v1 + v2
            par = payment / 2

            value = v7
            for j in range(1, 121):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7101 = v7 + par * 120
            v101 = value - v7101

            value = v7
            for j in range(1, 181):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7151 = v7 + par * 180
            v151 = value - v7151

            value = v7
            for j in range(1, 241):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7201 = v7 + par * 240
            v201 = value - v7201

            ### criando o segundo gráfico de 10 anos da simulação com metade da contribuição
            ax = pd.DataFrame([[v7101, v101]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7101, symbol=True, grouping=True,
                                          international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v101, symbol=True, grouping=True,
                                          international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right',
                      fontsize='x-small', frameon=False)
            ax.set_ylim(top=v7101 + v201)
            plt.savefig(f'{saida}foto4.png', transparent=True)
            plt.clf()
            plt.close()

            ### criando o segundo gráfico de 15 anos da simulação com metade da contribuição
            ax = pd.DataFrame([[v7151, v151]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7151, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v151, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right',
                      fontsize='x-small', frameon=False)
            ax.set_ylim(top=v7151 + v201)
            plt.savefig(f'{saida}foto5.png', transparent=True)
            plt.clf()
            plt.close()

            ### criando o segundo gráfico de 20 anos da simulação com metade da contribuição
            ax = pd.DataFrame([[v7201, v201]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7201, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v201, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right',
                      fontsize='x-small', frameon=False)
            ax.set_ylim(top=v7201 + v201)
            plt.savefig(f'{saida}foto6.png', transparent=True)
            plt.clf()
            plt.close()
            time.sleep(1)

            # criando as labels da segunda simulação
            dez1 = locale.currency(v7101 + v101, symbol=True, grouping=True, international=False)
            quinze1 = locale.currency(v7151 + v151, symbol=True, grouping=True, international=False)
            vinte1 = locale.currency(v7201 + v201, symbol=True, grouping=True, international=False)

            ### calculando os juros compostos com a contribuição total
            par = payment
            value = v7
            for j in range(1, 121):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7102 = v7 + par * 120
            v102 = value - v7102

            value = v7
            for j in range(1, 181):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7152 = v7 + par * 180
            v152 = value - v7152

            value = v7
            for j in range(1, 241):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7202 = v7 + par * 240
            v202 = value - v7202

            ### criando o terceiro gráfico de 10 anos da simulação com a contribuição total
            ax = pd.DataFrame([[v7102, v102]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7102, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v102, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right',
                      fontsize='x-small', frameon=False)
            ax.set_ylim(top=v7102 + v202)
            plt.savefig(f'{saida}foto7.png', transparent=True)
            plt.clf()
            plt.close()

            ### criando o terceiro gráfico de 15 anos da simulação com a contribuição total
            ax = pd.DataFrame([[v7152, v152]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7152, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v152, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right',
                      fontsize='x-small', frameon=False)
            ax.set_ylim(top=v7152 + v202)
            plt.savefig(f'{saida}foto8.png', transparent=True)
            plt.clf()
            plt.close()

            ### criando o terceiro gráfico de 20 anos da simulação com a contribuição total
            ax = pd.DataFrame([[v7202, v202]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7202, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v202, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right',
                      fontsize='x-small', frameon=False)
            ax.set_ylim(top=v7202 + v202)
            plt.savefig(f'{saida}foto9.png', transparent=True)
            plt.clf()
            plt.close()
            time.sleep(1)

            # criando as labels da terceira simulação
            dez2 = locale.currency(v7102 + v102, symbol=True, grouping=True, international=False)
            quinze2 = locale.currency(v7152 + v152, symbol=True, grouping=True, international=False)
            vinte2 = locale.currency(v7202 + v202, symbol=True, grouping=True, international=False)

            # Chamando a função que troca as variáveis no PPT
            self.principal(cpf, valor, parcela, com, sem, renda, patrocina, percentual, resgate, bruto, id, template,
                           saida, aposentado, plano, dez, quinze, vinte, taxa, tag, pdf, dez1, quinze1, vinte1, dez2,
                           quinze2, vinte2, meia, cinquenta, text, payment, v10, v101, v102)

            # Contador
            self.quant += 1
            self.quantidade.set(f"Quantidade: {self.quant}")

    def valor_parcela(self, psa):
        path = f'{os.path.abspath("edgedriver")}'
        previous_version = open(os.path.join(path, "edgeversion.txt")).read().strip()
        installed_version = \
            os.popen("reg query HKCU\\Software\\Microsoft\\Edge\\BLBeacon /v version").read().split()[-1]
        if previous_version == installed_version:
            print('ja att')
        else:
            print('precisa att')
            with open(os.path.join(path, "edgeversion.txt"), "w") as version_file:
                version_file.write(installed_version)
            os.remove(os.path.join(path, "msedgedriver.exe"))
            download = f"https://msedgedriver.azureedge.net/{installed_version}/edgedriver_win64.zip"
            response = requests.get(download)
            zipp = "edgedriver//edgedriver.zip"
            with open(zipp, "wb") as f:
                f.write(response.content)
            with ZipFile(zipp, "r") as zip_ref:
                zip_ref.extractall("edgedriver")
            os.remove(zipp)

        direcao = f'{os.path.abspath("edgedriver")}\\msedgedriver.exe'
        options = EdgeOptions()
        options.use_chromium = True
        #options.add_argument("headless")
        options.add_argument("disable-gpu")
        directory = f'{os.path.abspath("template-resgate")}\\'
        if not os.path.exists(directory):
            os.makedirs(directory)
        options.add_experimental_option("prefs", {"download.default_directory": directory})
        browser = Edge(executable_path=direcao, options=options)
        browser.get('https://ww1.sinqiaprevidencia.com.br/')

        sleep(2)
        username = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/div[3]/input')
        username.send_keys(self.user.get())
        password = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/div[4]/input')
        password.send_keys(self.password.get())
        selector = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/select')
        drop = Select(selector)
        drop.select_by_visible_text("Previdência")
        sleep(1)
        enter = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/table/tbody/tr[3]/td[1]/input')
        enter.click()

        try:
            incorreto = '/html/body/div[1]/div[1]/form[1]/center/div'
            acesso = browser.find_element(By.XPATH, incorreto)
            self.primeiro_aviso = acesso.text
            self.segundo_aviso = ''
            self.aviso()
            self.stop()
            browser.quit()

        except:
            iframe = browser.find_element(By.CSS_SELECTOR, '#Iframe1')
            browser.switch_to.frame(iframe)
            lupa = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
            "/html/body/table/tbody/tr/td[4]/table/tbody/tr/td[14]/a/img")))
            lupa.click()
            sleep(2)
            value = []
            for i in psa:
                browser.switch_to.default_content()
                iform = browser.find_element(By.CSS_SELECTOR, '#iform')
                browser.switch_to.frame(iform)
                campo = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                '/html/body/form/table/tbody/tr[3]/td/table/tbody/tr/td[2]/table/tbody/tr/td/table/tbody/tr[3]/td[6]/input')))
                campo.clear()
                campo.send_keys(i)

                search = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                '/html/body/form/table/tbody/tr[3]/td/table/tbody/tr/td[3]/a/input')))
                search.click()

                browser.switch_to.default_content()

                ibusca = browser.find_element(By.CSS_SELECTOR, '#ibusca')
                browser.switch_to.frame(ibusca)
                participante = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                '/html/body/div[1]/table/tbody/tr/td[1]/table/tbody/tr/td[1]/a[2]')))
                participante.click()

                browser.switch_to.default_content()

                idetalhe = browser.find_element(By.CSS_SELECTOR, '#iDetalhe')
                browser.switch_to.frame(idetalhe)
                arrecadacao = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                '/html/body/table/tbody/tr/td/table/tbody/tr[1]/td/table/tbody/tr/td[3]/a/img')))
                arrecadacao.click()

                sit_atual = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                '/html/body/table/tbody/tr/td/table/tbody/tr[1]/td/table/tbody/tr/td[8]/table/tbody/tr[1]/td/table/tbody/tr[1]/td[8]/a/img')))
                sit_atual.click()

                tamanho = len(os.listdir(directory))
                browser.switch_to.default_content()
                idetalhe = browser.find_element(By.CSS_SELECTOR, '#iDetalhe')
                browser.switch_to.frame(idetalhe)
                ipcinco = browser.find_element(By.CSS_SELECTOR, '#p5')
                browser.switch_to.frame(ipcinco)
                ipquatroa1 = browser.find_element(By.CSS_SELECTOR, '#p4a')
                browser.switch_to.frame(ipquatroa1)
                download = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                '/html/body/table/tbody/tr[2]/td/table/tbody/tr/td[2]/div[2]/table/tbody/tr/td[2]/a/img')))
                download.click()
                try:
                    while True:
                        arq = os.listdir(directory)
                        arquivos = len(list((pal for pal in arq if '.xlsx' or '.xls' in pal)))
                        if tamanho < arquivos:
                            sleep(2)
                            files = os.listdir(directory)
                            files.sort(key=lambda x: os.path.getctime(os.path.join(directory, x)))
                            last_file = files[-1]
                            original = directory + '\\' + last_file
                            new = directory + 'parcela.xlsx'
                            os.rename(original, new)
                            break
                    data = pd.read_excel(directory + 'parcela.xlsx')
                    listra = []
                    for v in list(data[7:]['Unnamed: 1']):
                        if v != 0:
                            listra.append(v)
                    value.append(listra[1])
                    os.remove(directory + 'parcela.xlsx')
                except Exception:
                    value.append(0.00)
                browser.switch_to.default_content()
                sleep(1)
            browser.quit()
            return value

    def aviso(self):
        # Janela que gera os avisos
        aviso_janela = tk.Toplevel()
        p = PhotoImage(file='Base//logo.png')

        # Janela
        aviso_janela.iconphoto(False, p)
        aviso_janela.title("Gerador")
        aviso_janela.config(width=300, height=200, bg='#ffffff')
        aviso_janela.resizable(width=False, height=False)

        # Botão
        botao_aviso = tk.Button(aviso_janela, text="Fechar", command=aviso_janela.destroy)
        botao_aviso.config(bg='#ffffff', activebackground="#e6e6e6", activeforeground="Black", width=10)
        botao_aviso.place(x=105, y=150)

        # Label
        label_aviso = tk.Label(aviso_janela, text=str(self.primeiro_aviso) + '\n' + str(self.segundo_aviso))
        label_aviso.config(font=("Courier", 10), bg='#ffffff')
        label_aviso.place(x=50, y=60)


Main()
