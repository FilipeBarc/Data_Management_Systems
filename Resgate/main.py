# pip install pandas
# pip install openpyxl
# pip install fsspec
# pip install Pyarrow
# pip install matplotlib
# pip install python-pptx
# pip install comtypes
# pip install PyPDF2
# pip install pywin32
# pip install requests
# pip install pyinstaller
# pip install selenium
# pip install msedge-selenium-tools selenium==3.141

import tkinter as tk
from tkinter import *
import os
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import comtypes.client
import locale
import time
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Pt, Inches
from dateutil.relativedelta import relativedelta
from PyPDF2 import PdfWriter, PdfReader
import threading
import pythoncom
import requests
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from msedge.selenium_tools import Edge, EdgeOptions
from selenium.webdriver.support.ui import Select, WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from time import sleep
from zipfile import ZipFile
locale.setlocale(locale.LC_ALL, 'pt_BR.UTF-8')


class Main:

    def __init__(self):
        self.root = tk.Tk()
        self.root.config(width=230, height=520)
        self.root.resizable(width=False, height=False)
        p = PhotoImage(file='Base//logo.png')
        self.root.iconphoto(False, p)
        self.root.title('Resgate')
        self.root.config(bg='#ffffff')
        self.hoje = datetime.today()
        pythoncom.CoInitialize()

        # instancias definidas fora de __init__
        self.df = None
        self.parar = None
        self.primeiro_aviso = None
        self.segundo_aviso = None
        self.quant = None

        # Caixas
        self.caixa = tk.LabelFrame(self.root, text="Reporte", bd=5, width=20, height=20)
        self.caixa.config(bg='#ffffff')
        self.caixa.place(x=55, y=110)

        # Labels
        label = tk.Label(self.root, text="Gerador de informativo\npara retenção de resgate")
        label.config(font=("Arial", 10), bg='#ffffff')
        label.place(x=40, y=55)
        user_label = tk.Label(self.root, text="Usuário")
        user_label.config(font=("Arial", 10), bg='#ffffff')
        user_label.place(x=50, y=293)
        senha_label = tk.Label(self.root, text="Senha")
        senha_label.config(font=("Arial", 10), bg='#ffffff')
        senha_label.place(x=50, y=338)
        taxa = tk.Label(self.root, text="Taxa de juros anual (%)")
        taxa.config(font=("Arial", 10), bg='#ffffff')
        taxa.place(x=45, y=243)
        self.quantidade = tk.StringVar()
        self.quantidade.set("Quantidade: 0")
        self.label1 = tk.Label(self.caixa, textvariable=self.quantidade)
        self.label1.config(font=("Arial", 10), bg='#ffffff')
        self.label1.grid(row=0)
        self.counter = datetime(2000, 1, 1, 0, 0, 0)
        self.string = self.counter.strftime("Tempo: " + "%H:%M:%S")
        self.tempo = tk.Label(self.caixa)
        self.tempo.config(font=("Arial", 10), bg='#ffffff', text=self.string)
        self.tempo.grid(row=1, sticky="W")

        # Caixas de texto
        self.user = StringVar()
        self.login = tk.Entry(self.root, bg="#f2f2f2", textvariable=self.user)
        self.login.place(x=50, y=315)
        self.password = StringVar()
        self.senha = tk.Entry(self.root, bg="#f2f2f2", textvariable=self.password)
        self.senha.config(show="*")
        self.senha.place(x=50, y=360)
        self.taxa = StringVar()

        # SpinBox
        self.juros_spin = tk.Spinbox(self.root, from_=0, to=20, width=10, increment=0.5)
        self.juros_spin.delete(0, "end")
        self.juros_spin.insert(0, '8,5')
        self.juros_spin.place(x=50, y=265)

        # Botões
        self.botao_gerar = tk.Button(self.root, text="Gerar Template", command=self.funcao)
        self.botao_gerar.config(width=16, bg='#ffffff', activebackground="#e6e6e6", activeforeground="Black")
        self.botao_gerar.place(x=50, y=400)
        self.botao_reiniciar = tk.Button(self.root, text="Reiniciar", command=self.ativar)
        self.botao_reiniciar.config(width=16, bg='#ffffff', activebackground="#e6e6e6", activeforeground="Black")
        self.botao_reiniciar.place(x=50, y=430)
        self.botao_sair = tk.Button(self.root, text="Sair", command=self.root.destroy)
        self.botao_sair.config(width=16, bg='#ffffff', activebackground="#e6e6e6", activeforeground="Black")
        self.botao_sair.place(x=50, y=460)

        self.root.mainloop()

    def time(self):
        # Faz o contador rodar
        self.string = self.counter.strftime("Tempo: " + "%H:%M:%S")
        self.tempo.config(text=self.string)
        self.parar = self.root.after(1000, self.time)
        self.counter += timedelta(seconds=1)

    def funcao(self):
        # Verificar se o programa esta rodando pela primeira vez e pede para reiniciar
        if self.quantidade.get() == "Quantidade: 0":
            self.senha.config(state='disabled')
            self.login.config(state='disabled')
            self.juros_spin.config(state='disabled')
            self.time()
            threading.Thread(target=self.contador).start()
        else:
            self.primeiro_aviso = 'Clique no botão Reiniciar'
            self.segundo_aviso = ' para gerar novamente'
            self.aviso()

    def contador(self):
        # Zera o contador e inicia a escolha do caminho
        self.quant = 0
        self.editar()

    def stop(self):
        # Finaliza o programa
        self.root.after_cancel(self.parar)

    def ativar(self):
        # Reinicia o programa depois de rodar a primeira vez
        self.quantidade.set("Quantidade: 0")
        self.label1.config(text=self.quantidade)
        self.counter = datetime(2000, 1, 1, 0, 0, 0)
        self.string = self.counter.strftime("Tempo: " + "%H:%M:%S")
        self.tempo.config(text=self.string)
        self.senha.config(state='normal')
        self.login.config(state='normal')
        self.juros_spin.config(state='normal')

    def editar(self):
        # Manipulação da tabela para criação dos templates
        df = pd.read_excel('Base//tabela_modelo.xlsx')
        total = []
        for i in df.index:
            total.append(df['Saldo participante'][i] + df['Saldo patrocinadora'][i])
        df['Total'] = total

        dias = []
        if datetime.today().month == 1:
            anterior = f'25/12/{datetime.today().year - 1}'
        else:
            anterior = f'25/{str(datetime.today().month - 1).zfill(2)}/{datetime.today().year}'

        for i in df.index:
            novo = relativedelta(pd.to_datetime(anterior, format='%d/%m/%Y'),
                                 pd.to_datetime(df['DataAdmissao'][i], format='%d/%m/%Y'))
            dias.append(novo.years * 12 + novo.months)
        df['Dias em Meses'] = dias

        perc = []
        for i in df.index:
            if df['Plano'][i] == 'Visão Multi':
                perc.append(60 if df['Dias em Meses'][i] * 0.25 >= 60 else df['Dias em Meses'][i] * 0.25)
            else:
                perc.append(3 if df['Dias em Meses'][i] <= 12 else (
                    6 if df['Dias em Meses'][i] > 12 and df['Dias em Meses'][i] <= 24 else (
                    9 if df['Dias em Meses'][i] > 24 and df['Dias em Meses'][i] <= 36 else (
                    12 if df['Dias em Meses'][i] > 36 and df['Dias em Meses'][i] <= 48 else (
                    60 if df['Dias em Meses'][i] > 48 and df['Dias em Meses'][i] <= 60 else (
                    67.5 if df['Dias em Meses'][i] > 60 and df['Dias em Meses'][i] <= 72 else (
                    75 if df['Dias em Meses'][i] > 72 and df['Dias em Meses'][i] <= 84 else (
                    82.5 if df['Dias em Meses'][i] > 84 and df['Dias em Meses'][i] <= 96 else (
                    90 if df['Dias em Meses'][i] > 96 else df['Dias em Meses'][i])))))))))
        df['Percentual de resgate'] = perc

        val = []
        for i in df.index:
            val.append(df['Saldo patrocinadora'][i] * df['Percentual de resgate'][i] / 100)
        df['Valor resgatavel'] = val

        bruto = []
        for i in df.index:
            bruto.append(df['Saldo participante'][i] + df['Valor resgatavel'][i])
        df['Valor bruto de resgate'] = bruto

        nomes = []
        for i in df['Nome'].index:
            nomes.append(df['Nome'][i].split(' ')[0].capitalize())
        df['Primeiro nome'] = nomes

        lista_sa = list(df['ParticipanteSA'])
        payment = self.valor_parcela(lista_sa)
        df['Pagamentos'] = payment

        df_reg = self.regressivo(lista_sa)

        df = df[['CPF', 'ParticipanteSA', 'Nome', 'Primeiro nome', 'Plano', 'Saldo participante', 'Saldo patrocinadora',
                 'Total', 'Percentual de resgate', 'Valor resgatavel', 'Valor bruto de resgate', 'Dias em Meses',
                 'Pagamentos']]

        df = pd.concat([df, df_reg], axis=1)

        #df.to_excel("Templates/tabela final.xlsx")

        self.gerador(df)
        self.primeiro_aviso = 'Comunicados gerados com'
        self.segundo_aviso = ' sucesso!'
        self.aviso()
        self.stop()

    # Função que garante a formatação original dos textos
    def manter_formatacao_original(self, paragrafo, fonte_original):
        for run in paragrafo.runs:
            run.font.size = fonte_original.size
            run.font.bold = fonte_original.bold
            run.font.italic = fonte_original.italic
            run.font.color.rgb = fonte_original.color.rgb
            run.font.name = fonte_original.name
            run.font.underline = fonte_original.underline

    # Função que faz a substituição dos textos
    def substituir_texto(self, apresentacao, texto_antigo, texto_novo):
        for slide in apresentacao.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    for paragrafo in shape.text_frame.paragraphs:
                        for run in paragrafo.runs:
                            if texto_antigo in run.text:
                                fonte_original = run.font
                                run.text = run.text.replace(texto_antigo, texto_novo)
                                novo_paragrafo = shape.text_frame.add_paragraph()
                                self.manter_formatacao_original(novo_paragrafo, fonte_original)

    # Função que busca as chaves dentro o PPT
    def principal(self, cpf, nome, valor, parcela, renda, patrocina, percentual, resgate, bruto, id, ppt, saida, plano,
                  dez, quinze, vinte, imposto, liquido, taxa, pdf, dez1, quinze1, vinte1, dez2, quinze2, vinte2, meia,
                  payment, saldo10, saldo15, saldo20, saldo25, saldo30, saldo35, saldo40, imposto10, imposto15,
                  imposto20, imposto25, imposto30, imposto35, imposto40, liquido10, liquido15, liquido20, liquido25,
                  liquido30, liquido35, liquido40, v10, v15, v20):

        apresentacao = Presentation(ppt)
        ppt_out = f'{saida}{id}.pptx'

        self.substituir_texto(apresentacao, '{nome}', nome)
        self.substituir_texto(apresentacao, '{renda}', renda)
        self.substituir_texto(apresentacao, '{patrocinadora}', patrocina)
        self.substituir_texto(apresentacao, '{valor}', valor)
        self.substituir_texto(apresentacao, '{percentual}', percentual)
        self.substituir_texto(apresentacao, '{resgate}', resgate)
        self.substituir_texto(apresentacao, '{bruto}', bruto)
        self.substituir_texto(apresentacao, '{plano}', plano)
        self.substituir_texto(apresentacao, '{imposto}', imposto)
        self.substituir_texto(apresentacao, '{liquido}', liquido)
        self.substituir_texto(apresentacao, '{taxa}', taxa)
        self.substituir_texto(apresentacao, '{parcela}', parcela)
        self.substituir_texto(apresentacao, '{meia}', meia)

        self.substituir_texto(apresentacao, '{dez}', dez)
        self.substituir_texto(apresentacao, '{quinze}', quinze)
        self.substituir_texto(apresentacao, '{vinte}', vinte)
        self.substituir_texto(apresentacao, '{dez1}', dez1)
        self.substituir_texto(apresentacao, '{quinze1}', quinze1)
        self.substituir_texto(apresentacao, '{vinte1}', vinte1)
        self.substituir_texto(apresentacao, '{dez2}', dez2)
        self.substituir_texto(apresentacao, '{quinze2}', quinze2)
        self.substituir_texto(apresentacao, '{vinte2}', vinte2)

        self.substituir_texto(apresentacao, '{saldo10}', saldo10)
        self.substituir_texto(apresentacao, '{saldo15}', saldo15)
        self.substituir_texto(apresentacao, '{saldo20}', saldo20)
        self.substituir_texto(apresentacao, '{saldo25}', saldo25)
        self.substituir_texto(apresentacao, '{saldo30}', saldo30)
        self.substituir_texto(apresentacao, '{saldo35}', saldo35)
        self.substituir_texto(apresentacao, '{saldo40}', saldo40)

        self.substituir_texto(apresentacao, '{imposto10}', imposto10)
        self.substituir_texto(apresentacao, '{imposto15}', imposto15)
        self.substituir_texto(apresentacao, '{imposto20}', imposto20)
        self.substituir_texto(apresentacao, '{imposto25}', imposto25)
        self.substituir_texto(apresentacao, '{imposto30}', imposto30)
        self.substituir_texto(apresentacao, '{imposto35}', imposto35)
        self.substituir_texto(apresentacao, '{imposto40}', imposto40)

        self.substituir_texto(apresentacao, '{liquido10}', liquido10)
        self.substituir_texto(apresentacao, '{liquido15}', liquido15)
        self.substituir_texto(apresentacao, '{liquido20}', liquido20)
        self.substituir_texto(apresentacao, '{liquido25}', liquido25)
        self.substituir_texto(apresentacao, '{liquido30}', liquido30)
        self.substituir_texto(apresentacao, '{liquido35}', liquido35)
        self.substituir_texto(apresentacao, '{liquido40}', liquido40)

        apresentacao.save(ppt_out)
        self.imagens(ppt_out, cpf, saida, pdf, id, payment, v10, v15, v20)

    # Função que insere as imagens no PPT
    def imagens(self, ppt_out, cpf, saida, pdf, id, payment, v10, v15, v20):
        prs = Presentation(ppt_out)
        slide = prs.slides[1]
        img_path = f"{saida}foto.png"
        left = Inches(6.5)
        top = Inches(1.5)
        width = Inches(13.5)
        height = Inches(10)
        slide.shapes.add_picture(img_path, left, top, width, height)
        os.remove(f"{saida}foto.png")

        slide = prs.slides[2]
        img_path1 = f"{saida}foto1.png"
        left1 = Inches(0.2)
        top1 = Inches(4.3)
        width1 = Inches(8.2)
        height1 = Inches(6)
        img_path2 = f"{saida}foto2.png"
        left2 = Inches(5.8)
        img_path3 = f"{saida}foto3.png"
        left3 = Inches(11.5)
        slide.shapes.add_picture(img_path1, left1, top1, width1, height1)
        slide.shapes.add_picture(img_path2, left2, top1, width1, height1)
        slide.shapes.add_picture(img_path3, left3, top1, width1, height1)
        slide = prs.slides[3]
        img_path1 = f"{saida}foto4.png"
        left1 = Inches(0.2)
        top1 = Inches(4.3)
        width1 = Inches(8.2)
        height1 = Inches(6)
        img_path2 = f"{saida}foto5.png"
        left2 = Inches(5.8)
        img_path3 = f"{saida}foto6.png"
        left3 = Inches(11.5)
        slide.shapes.add_picture(img_path1, left1, top1, width1, height1)
        slide.shapes.add_picture(img_path2, left2, top1, width1, height1)
        slide.shapes.add_picture(img_path3, left3, top1, width1, height1)
        slide = prs.slides[4]
        img_path1 = f"{saida}foto7.png"
        left1 = Inches(0.2)
        top1 = Inches(4.3)
        width1 = Inches(8.2)
        height1 = Inches(6)
        img_path2 = f"{saida}foto8.png"
        left2 = Inches(5.8)
        img_path3 = f"{saida}foto9.png"
        left3 = Inches(11.5)
        slide.shapes.add_picture(img_path1, left1, top1, width1, height1)
        slide.shapes.add_picture(img_path2, left2, top1, width1, height1)
        slide.shapes.add_picture(img_path3, left3, top1, width1, height1)

        if v10 <= 0 or v15 <= 0 or v20 <= 0:
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[2])

        if payment == 0:
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[3])
            xml_slides.remove(slides[4])

        os.remove(f"{saida}foto1.png")
        os.remove(f"{saida}foto2.png")
        os.remove(f"{saida}foto3.png")
        os.remove(f"{saida}foto4.png")
        os.remove(f"{saida}foto5.png")
        os.remove(f"{saida}foto6.png")
        os.remove(f"{saida}foto7.png")
        os.remove(f"{saida}foto8.png")
        os.remove(f"{saida}foto9.png")

        prs.save(f'{saida}{id}.pptx')
        time.sleep(1)
        self.ppttopdf(cpf, pdf, id, saida)

    # Função que transforma o PPT em PDF
    def ppttopdf(self, cpf, pdf, id, out):
        entrada = f'{pdf}{id}.pptx'
        saida = f'{pdf}{id}.pdf'
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application", pythoncom.CoInitialize())
        deck = powerpoint.Presentations.Open(entrada)
        deck.ExportAsFixedFormat(saida, 32)
        deck.Close()
        powerpoint.Quit()
        sleep(1)
        self.encriptar(cpf, saida, entrada)

    # Função que faz a encriptação do PDF
    def encriptar(self, cpf, saida, entrada):
        out = PdfWriter()
        file = PdfReader(saida)
        num = len(file.pages)
        os.remove(entrada)
        for idx in range(num):
            page = file.pages[idx]
            out.add_page(page)
        password = cpf[-4:]
        out.encrypt(password)
        with open(saida, "wb") as f:
            out.write(f)

    # Função que coleta as variáveis da tabela e cria os gráficos das 3 diferentes opções para inserir no PPT
    def gerador(self, base):
        template = 'Base//Retenção de resgate.pptx'
        pdf = f'{os.path.abspath("Templates")}\\'
        saida = "Templates//"
        # Primeira opção dos resgates
        for i in base.index:
            # Coleta das variáveis
            cpf = str(base['CPF'][i])
            sa = str(base['ParticipanteSA'][i])
            nome = base['Primeiro nome'][i]
            renda = locale.currency(base['Saldo participante'][i], symbol=True, grouping=True, international=False)
            patrocina = locale.currency(base['Saldo patrocinadora'][i], symbol=True, grouping=True, international=False)
            valor = locale.currency(base['Total'][i], symbol=True, grouping=True, international=False)
            percentual = str(base['Percentual de resgate'][i]) + '%'
            resgate = locale.currency(base['Valor resgatavel'][i], symbol=True, grouping=True, international=False)
            bruto = locale.currency(base['Valor bruto de resgate'][i], symbol=True, grouping=True, international=False)
            plano = base['Plano'][i]
            payment = base['Pagamentos'][i]
            parcela = locale.currency(payment, symbol=True, grouping=True, international=False)
            meiaparcela = locale.currency(payment / 2, symbol=True, grouping=True, international=False)

            sal10 = base['saldo10'][i]
            sal15 = base['saldo15'][i]
            sal20 = base['saldo20'][i]
            sal25 = base['saldo25'][i]
            sal30 = base['saldo30'][i]
            sal35 = base['saldo35'][i]
            sal40 = base['saldo40'][i]

            saldo10 = locale.currency(sal10, symbol=True, grouping=True, international=False)
            saldo15 = locale.currency(sal15, symbol=True, grouping=True, international=False)
            saldo20 = locale.currency(sal20, symbol=True, grouping=True, international=False)
            saldo25 = locale.currency(sal25, symbol=True, grouping=True, international=False)
            saldo30 = locale.currency(sal30, symbol=True, grouping=True, international=False)
            saldo35 = locale.currency(sal35, symbol=True, grouping=True, international=False)
            saldo40 = locale.currency(sal40, symbol=True, grouping=True, international=False)

            imp40 = sal10 * 0.1 + sal15 * 0.15 + sal20 * 0.2 + sal25 * 0.25 + sal30 * 0.3 + sal35 * 0.35
            imposto10 = locale.currency(sal10 * 0.1, symbol=True, grouping=True, international=False)
            imposto15 = locale.currency(sal15 * 0.15, symbol=True, grouping=True, international=False)
            imposto20 = locale.currency(sal20 * 0.2, symbol=True, grouping=True, international=False)
            imposto25 = locale.currency(sal25 * 0.25, symbol=True, grouping=True, international=False)
            imposto30 = locale.currency(sal30 * 0.3, symbol=True, grouping=True, international=False)
            imposto35 = locale.currency(sal35 * 0.35, symbol=True, grouping=True, international=False)
            imposto40 = locale.currency(imp40, symbol=True, grouping=True, international=False)

            liq40 = ((sal10 - sal10 * 0.1) + (sal15 - sal15 * 0.15) + (sal20 - sal20 * 0.2) + (sal25 - sal25 * 0.25) +
                     (sal30 - sal30 * 0.3) + (sal35 - sal35 * 0.35))
            liquido10 = locale.currency(sal10 - sal10 * 0.1, symbol=True, grouping=True, international=False)
            liquido15 = locale.currency(sal15 - sal15 * 0.15, symbol=True, grouping=True, international=False)
            liquido20 = locale.currency(sal20 - sal20 * 0.2, symbol=True, grouping=True, international=False)
            liquido25 = locale.currency(sal25 - sal25 * 0.25, symbol=True, grouping=True, international=False)
            liquido30 = locale.currency(sal30 - sal30 * 0.3, symbol=True, grouping=True, international=False)
            liquido35 = locale.currency(sal35 - sal35 * 0.35, symbol=True, grouping=True, international=False)
            liquido40 = locale.currency(liq40, symbol=True, grouping=True, international=False)

            # Criação da primeira imagem em gráfico pizza
            fig = plt.figure()
            ax = fig.add_subplot(111)
            explode = (0.01, 0.01)
            v1 = list(base[['Saldo participante', 'Saldo patrocinadora']].iloc[i])[0]
            v2 = list(base[['Saldo participante', 'Saldo patrocinadora']].iloc[i])[1]
            ax = base[['Saldo participante', 'Saldo patrocinadora']].iloc[i].plot.pie(colors=["#58aa7a", "#065d70"],
                shadow=True, explode=explode, labels=None, autopct=lambda p: locale.currency(p * (v1 + v2) / 100,
                symbol=True, grouping=True, international=False), textprops={'color': "w", 'weight': 'bold'})
            ax.yaxis.set_visible(False)
            ax.legend(loc='lower right', fontsize='x-small', labels=['Participante', 'Patrocinadora'],
                      frameon=False)
            plt.savefig(f'{saida}foto.png', transparent=True)
            plt.close()

            # Calculo do juros compostos para os 3 gráficos do terceiro slide
            v3 = v1 + v2
            taxa = str(self.juros_spin.get())
            tx = float(taxa.replace(',','.')) / 100
            tx_mes = (1 + tx) ** (1 / 12) - 1

            if plano == 'Visão Multi':
                pri = 55.00 / 12
                seg = (1 + 0.0069) ** (1 / 12) - 1
                ter = (1 + 0.0046) ** (1 / 12) - 1
                qua = (1 + 0.0028) ** (1 / 12) - 1
                qui = (1 + 0.0014) ** (1 / 12) - 1
                sex = 1900.00 / 12
            else:
                pri = 93.00 / 12
                seg = (1 + 0.0102) ** (1 / 12) - 1
                ter = (1 + 0.0092) ** (1 / 12) - 1
                qua = (1 + 0.0077) ** (1 / 12) - 1
                qui = (1 + 0.0037) ** (1 / 12) - 1
                sex = 3140.00 / 12

            value = v3
            for j in range(1, 121):
                if value <= 10000:
                    value = value + value * tx_mes - pri
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1)
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1)
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1)
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1)
                elif value > 1250000:
                    value = value + value * tx_mes - sex
            v10 = value - v3

            value = v3
            for j in range(1, 181):
                if value <= 10000:
                    value = value + value * tx_mes - pri
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1)
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1)
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1)
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1)
                elif value > 1250000:
                    value = value + value * tx_mes - sex
            v15 = value - v3

            value = v3
            for j in range(1, 241):
                if value <= 10000:
                    value = value + value * tx_mes - pri
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1)
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1)
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1)
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1)
                elif value > 1250000:
                    value = value + value * tx_mes - sex
            v20 = value - v3

            # Criação do primeiro gráfico de dez anos
            ax = pd.DataFrame([[v3, v10]],columns=['Contribuição','Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v3, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v10, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right', fontsize='x-small',
                      frameon=False)
            ax.set_ylim(top=v3 + v20)
            plt.savefig(f'{saida}foto1.png', transparent=True)
            plt.close()

            # Criação do segundo gráfico de 15 anos
            ax = pd.DataFrame([[v3, v15]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v3, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v15, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right', fontsize='x-small',
                      frameon=False)
            ax.set_ylim(top=v3 + v20)
            plt.savefig(f'{saida}foto2.png', transparent=True)
            plt.close()

            # Criação do terceiro gráfico de 15 anos
            ax = pd.DataFrame([[v3, v20]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v3, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v20, symbol=True, grouping=True,
                                      international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right', fontsize='x-small',
                      frameon=False)
            ax.set_ylim(top=v3 + v20)
            plt.savefig(f'{saida}foto3.png', transparent=True)
            plt.close()

            # Criação das variáveis de 10, 15 e 20 anos, e as de valor bruto, imposto e valor líquido
            dez = locale.currency(v3 + v10, symbol=True, grouping=True, international=False)
            quinze = locale.currency(v3 + v15, symbol=True, grouping=True, international=False)
            vinte = locale.currency(v3 + v20, symbol=True, grouping=True, international=False)
            valor_bruto = base['Valor bruto de resgate'][i]
            imposto = locale.currency(valor_bruto * 0.15, symbol=True, grouping=True, international=False)
            liquido = locale.currency(valor_bruto - valor_bruto * 0.15, symbol=True, grouping=True, international=False)

            ### calculando os juros compostos com metade da contribuição
            v7 = v1 + v2
            par = payment / 2
            value = v7

            for j in range(1, 121):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7101 = v7 + par * 120
            v101 = value - v7101
            value = v7

            for j in range(1, 181):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7151 = v7 + par * 180
            v151 = value - v7151
            value = v7

            for j in range(1, 241):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7201 = v7 + par * 240
            v201 = value - v7201

            ### criando o segundo gráfico de 10 anos da simulação com metade da comtribuição
            ax = pd.DataFrame([[v7101, v101]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7101, symbol=True, grouping=True,
                                          international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v101, symbol=True, grouping=True,
                                          international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right', fontsize='x-small',
                      frameon=False)
            ax.set_ylim(top=v7101 + v201)
            plt.savefig(f'{saida}foto4.png', transparent=True)
            plt.close()

            ### criando o segundo gráfico de 15 anos da simulação com metade da comtribuição
            ax = pd.DataFrame([[v7151, v151]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7151, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v151, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right', fontsize='x-small',
                      frameon=False)
            ax.set_ylim(top=v7151 + v201)
            plt.savefig(f'{saida}foto5.png', transparent=True)
            plt.close()

            ### criando o segundo gráfico de 20 anos da simulação com metade da comtribuição
            ax = pd.DataFrame([[v7201, v201]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7201, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v201, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right', fontsize='x-small',
                      frameon=False)
            ax.set_ylim(top=v7201 + v201)
            plt.savefig(f'{saida}foto6.png', transparent=True)
            plt.close()
            time.sleep(1)

            # criando as labels da segunda simulação
            dez1 = locale.currency(v7101 + v101, symbol=True, grouping=True, international=False)
            quinze1 = locale.currency(v7151 + v151, symbol=True, grouping=True, international=False)
            vinte1 = locale.currency(v7201 + v201, symbol=True, grouping=True, international=False)

            ### calculando os juros compostos com a comtribuição total
            par = payment
            value = v7
            for j in range(1, 121):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7102 = v7 + par * 120
            v102 = value - v7102

            value = v7
            for j in range(1, 181):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7152 = v7 + par * 180
            v152 = value - v7152

            value = v7
            for j in range(1, 241):
                if value <= 10000:
                    value = value + value * tx_mes - pri + par
                elif value > 10000 and value <= 50000:
                    value = value + value * ((1 + tx_mes) / (1 + seg) - 1) + par
                elif value > 50000 and value <= 200000:
                    value = value + value * ((1 + tx_mes) / (1 + ter) - 1) + par
                elif value > 200000 and value <= 500000:
                    value = value + value * ((1 + tx_mes) / (1 + qua) - 1) + par
                elif value > 500000 and value <= 1250000:
                    value = value + value * ((1 + tx_mes) / (1 + qui) - 1) + par
                elif value > 1250000:
                    value = value + value * tx_mes - sex + par
            v7202 = v7 + par * 240
            v202 = value - v7202

            ### criando o terceiro gráfico de 10 anos da simulação com a comtribuição total
            ax = pd.DataFrame([[v7102, v102]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7102, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v102, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right', fontsize='x-small',
                      frameon=False)
            ax.set_ylim(top=v7102 + v202)
            plt.savefig(f'{saida}foto7.png', transparent=True)
            plt.close()

            ### criando o terceiro gráfico de 15 anos da simulação com a comtribuição total
            ax = pd.DataFrame([[v7152, v152]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7152, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v152, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right', fontsize='x-small',
                      frameon=False)
            ax.set_ylim(top=v7152 + v202)
            plt.savefig(f'{saida}foto8.png', transparent=True)
            plt.close()

            ### criando o terceiro gráfico de 20 anos da simulação com a comtribuição total
            ax = pd.DataFrame([[v7202, v202]], columns=['Contribuição', 'Rentabilidade']).plot.bar(
                color=["#58aa7a", "#065d70"], stacked=True, edgecolor="w")
            ax.bar_label(ax.containers[0], labels=[locale.currency(v7202, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.bar_label(ax.containers[1], labels=[locale.currency(v202, symbol=True, grouping=True,
                                           international=False)], label_type='center', color='w', weight='bold')
            ax.yaxis.set_visible(False)
            ax.xaxis.set_visible(False)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            ax.spines['left'].set_visible(False)
            handles, labels = plt.gca().get_legend_handles_labels()
            order = [1, 0]
            ax.legend([handles[i] for i in order], [labels[i] for i in order], loc='lower right', fontsize='x-small',
                      frameon=False)
            ax.set_ylim(top=v7202 + v202)
            plt.savefig(f'{saida}foto9.png', transparent=True)
            plt.close()
            time.sleep(1)

            # criando as labels da terceira simulação
            dez2 = locale.currency(v7102 + v102, symbol=True, grouping=True, international=False)
            quinze2 = locale.currency(v7152 + v152, symbol=True, grouping=True, international=False)
            vinte2 = locale.currency(v7202 + v202, symbol=True, grouping=True, international=False)

            # Chamando a função que troca as variáveis no PPT
            self.principal(cpf, nome, valor, parcela, renda, patrocina, percentual, resgate, bruto, sa, template,
                           saida, plano, dez, quinze, vinte, imposto, liquido, taxa, pdf, dez1, quinze1, vinte1,
                           dez2, quinze2, vinte2, meiaparcela, payment, saldo10, saldo15, saldo20, saldo25, saldo30,
                           saldo35, saldo40, imposto10, imposto15, imposto20, imposto25, imposto30, imposto35, imposto40,
                           liquido10, liquido15, liquido20, liquido25, liquido30, liquido35, liquido40, v10, v15, v20)

            # Contador
            self.quant += 1
            self.quantidade.set(f"Quantidade: {self.quant}")


    def valor_parcela(self, psa):
        path = f'{os.path.abspath("edgedriver")}'
        previous_version = open(os.path.join(path, "edgeversion.txt")).read().strip()
        installed_version = \
            os.popen("reg query HKCU\\Software\\Microsoft\\Edge\\BLBeacon /v version").read().split()[-1]
        if previous_version == installed_version:
            print('ja att')
        else:
            print('precisa att')
            with open(os.path.join(path, "edgeversion.txt"), "w") as version_file:
                version_file.write(installed_version)
            os.remove(os.path.join(path, "msedgedriver.exe"))
            download = f"https://msedgedriver.azureedge.net/{installed_version}/edgedriver_win64.zip"
            response = requests.get(download)
            zipp = "edgedriver//edgedriver.zip"
            with open(zipp, "wb") as f:
                f.write(response.content)
            with ZipFile(zipp, "r") as zip_ref:
                zip_ref.extractall("edgedriver")
            os.remove(zipp)

        direcao = f'{os.path.abspath("edgedriver")}\\msedgedriver.exe'
        options = EdgeOptions()
        options.use_chromium = True
        #options.add_argument("headless")
        options.add_argument("disable-gpu")
        directory = f'{os.path.abspath("Templates")}\\'
        if not os.path.exists(directory):
            os.makedirs(directory)
        options.add_experimental_option("prefs", {"download.default_directory": directory})
        browser = Edge(executable_path=direcao, options=options)
        browser.get('https://ww1.sinqiaprevidencia.com.br/')

        sleep(2)
        username = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/div[3]/input')
        username.send_keys(self.user.get())
        password = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/div[4]/input')
        password.send_keys(self.password.get())
        selector = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/select')
        drop = Select(selector)
        drop.select_by_visible_text("Previdência")
        sleep(1)
        enter = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/table/tbody/tr[3]/td[1]/input')
        enter.click()

        try:
            incorreto = '/html/body/div[1]/div[1]/form[1]/center/div'
            acesso = browser.find_element(By.XPATH, incorreto)
            self.primeiro_aviso = acesso.text
            self.segundo_aviso = ''
            self.aviso()
            self.stop()
            browser.quit()

        except:
            iframe = browser.find_element(By.CSS_SELECTOR, '#Iframe1')
            browser.switch_to.frame(iframe)
            lupa = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
            "/html/body/table/tbody/tr/td[4]/table/tbody/tr/td[14]/a/img")))
            lupa.click()

            sleep(2)
            value = []
            for i in psa:
                browser.switch_to.default_content()
                iform = browser.find_element(By.CSS_SELECTOR, '#iform')
                browser.switch_to.frame(iform)
                campo = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                        '/html/body/form/table/tbody/tr[3]/td/table/tbody/tr/td[2]/table/tbody/tr/td/table/tbody/tr[3]/td[6]/input')))
                campo.clear()
                campo.send_keys(i)
                search = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                         '/html/body/form/table/tbody/tr[3]/td/table/tbody/tr/td[3]/a/input')))
                search.click()
                browser.switch_to.default_content()
                ibusca = browser.find_element(By.CSS_SELECTOR, '#ibusca')
                browser.switch_to.frame(ibusca)
                sleep(2)
                try:
                    participante = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                    '/html/body/div[1]/table/tbody/tr/td[1]/table/tbody/tr/td[1]/a[2]')))
                    participante.click()
                except Exception:
                    self.primeiro_aviso = 'Erro no Participante SA'
                    self.segundo_aviso = f' {i}'
                    self.aviso()
                    self.stop()
                    browser.quit()
                browser.switch_to.default_content()
                idetalhe = browser.find_element(By.CSS_SELECTOR, '#iDetalhe')
                browser.switch_to.frame(idetalhe)
                arrecadacao = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                '/html/body/table/tbody/tr/td/table/tbody/tr[1]/td/table/tbody/tr/td[3]/a/img')))
                arrecadacao.click()
                sit_atual = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                '/html/body/table/tbody/tr/td/table/tbody/tr[1]/td/table/tbody/tr/td[8]/table/tbody/tr[1]/td/table/tbody/tr[1]/td[8]/a/img')))
                sit_atual.click()

                tamanho = len(os.listdir(directory))
                browser.switch_to.default_content()
                idetalhe = browser.find_element(By.CSS_SELECTOR, '#iDetalhe')
                browser.switch_to.frame(idetalhe)
                ipcinco = browser.find_element(By.CSS_SELECTOR, '#p5')
                browser.switch_to.frame(ipcinco)
                ipquatroa1 = browser.find_element(By.CSS_SELECTOR, '#p4a')
                browser.switch_to.frame(ipquatroa1)
                download = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                '/html/body/table/tbody/tr[2]/td/table/tbody/tr/td[2]/div[2]/table/tbody/tr/td[2]/a/img')))
                download.click()
                try:
                    while True:
                        arq = os.listdir(directory)
                        arquivos = len(list((pal for pal in arq if '.xlsx' or '.xls' in pal)))
                        if tamanho < arquivos:
                            sleep(2)
                            files = os.listdir(directory)
                            files.sort(key=lambda x: os.path.getctime(os.path.join(directory, x)))
                            last_file = files[-1]
                            original = directory + '\\' + last_file
                            new = directory + 'parcela.xlsx'
                            os.rename(original, new)
                            break
                    data = pd.read_excel(directory + 'parcela.xlsx')

                    listra = []
                    for v in list(data[7:]['Unnamed: 1']):
                        if v != 0:
                            listra.append(v)
                    try:
                        value.append(listra[1])
                    except Exception:
                        value.append(listra[0])
                    os.remove(directory + 'parcela.xlsx')
                except Exception:
                    value.append(0.00)
                browser.switch_to.default_content()
                sleep(1)
            browser.quit()
            return value

    def regressivo(self, lista):
        saldo10 = []
        saldo15 = []
        saldo20 = []
        saldo25 = []
        saldo30 = []
        saldo35 = []
        saldo40 = []

        path = f'{os.path.abspath("edgedriver")}'
        previous_version = open(os.path.join(path, "edgeversion.txt")).read().strip()
        installed_version = \
            os.popen("reg query HKCU\\Software\\Microsoft\\Edge\\BLBeacon /v version").read().split()[-1]

        if previous_version == installed_version:
            print('ja att')
        else:
            print('precisa att')
            with open(os.path.join(path, "edgeversion.txt"), "w") as version_file:
                version_file.write(installed_version)
            os.remove(os.path.join(path, "msedgedriver.exe"))
            download = f"https://msedgedriver.azureedge.net/{installed_version}/edgedriver_win64.zip"
            response = requests.get(download)
            zipp = "edgedriver//edgedriver.zip"
            with open(zipp, "wb") as f:
                f.write(response.content)
            with ZipFile(zipp, "r") as zip_ref:
                zip_ref.extractall("edgedriver")
            os.remove(zipp)

        direcao = f'{os.path.abspath("edgedriver")}\\msedgedriver.exe'
        options = EdgeOptions()
        options.use_chromium = True
        #options.add_argument("headless")
        options.add_argument("disable-gpu")
        directory = f'{os.path.abspath("Templates")}\\'
        if not os.path.exists(directory):
            os.makedirs(directory)
        options.add_experimental_option("prefs", {"download.default_directory": directory})
        browser = Edge(executable_path=direcao, options=options)
        browser.get('https://ww1.sinqiaprevidencia.com.br/')

        sleep(2)
        username = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/div[3]/input')
        username.send_keys(self.user.get())
        password = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/div[4]/input')
        password.send_keys(self.password.get())
        selector = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/select')
        drop = Select(selector)
        drop.select_by_visible_text("Previdência")
        sleep(1)
        enter = browser.find_element(By.XPATH, '/html/body/div[1]/div[1]/form[1]/table/tbody/tr[3]/td[1]/input')
        enter.click()

        try:
            incorreto = '/html/body/div[1]/div[1]/form[1]/center/div'
            acesso = browser.find_element(By.XPATH, incorreto)
            self.primeiro_aviso = acesso.text
            self.segundo_aviso = ''
            self.aviso()
            self.stop()
            browser.quit()

        except:
            iframe = browser.find_element(By.CSS_SELECTOR, '#Iframe1')
            browser.switch_to.frame(iframe)
            lupa = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
            "/html/body/table/tbody/tr/td[4]/table/tbody/tr/td[14]/a/img")))
            lupa.click()

        sleep(2)
        for i in lista:
            try:
                browser.switch_to.default_content()
                iform = browser.find_element(By.CSS_SELECTOR, '#iform')
                browser.switch_to.frame(iform)
                psa = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                      '/html/body/form/table/tbody/tr[3]/td/table/tbody/tr/td[2]/table/tbody/tr/td/table/tbody/tr[3]/td[6]/input')))
                psa.clear()
                psa.send_keys(i)
                search = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                         '/html/body/form/table/tbody/tr[3]/td/table/tbody/tr/td[3]/a/input')))
                search.click()
                browser.switch_to.default_content()
                ibusca = browser.find_element(By.CSS_SELECTOR, '#ibusca')
                browser.switch_to.frame(ibusca)
                participante = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                               '/html/body/div[1]/table/tbody/tr/td[1]/table/tbody/tr/td[1]/a[2]')))
                participante.click()
                browser.switch_to.default_content()
                idetalhe = browser.find_element(By.CSS_SELECTOR, '#iDetalhe')
                browser.switch_to.frame(idetalhe)
                cadastro = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                           '/html/body/table/tbody/tr/td/table/tbody/tr[1]/td/table/tbody/tr/td[2]/a/img')))
                cadastro.click()
                ipcinco = browser.find_element(By.CSS_SELECTOR, '#p5')
                browser.switch_to.frame(ipcinco)
                ipquatroa1 = browser.find_element(By.CSS_SELECTOR, '#p4a')
                browser.switch_to.frame(ipquatroa1)
                descer = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                         '/html/body/table/tbody/tr[1]/td/table/tbody/tr/td[4]/table/tbody/tr/td[2]/a/img')))
                descer.click()
                sleep(4)
                plano = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                        '/html/body/div/table/tbody/tr/td[1]/table/tbody/tr[1]/td[2]/table/tbody/tr[1]/td/table/tbody/tr[3]/td[2]/table/tbody/tr[2]/td[2]')))
                plano_nome = plano.text
                data_admissao = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                                '/html/body/div/table/tbody/tr/td[1]/table/tbody/tr[1]/td[2]/table/tbody/tr[1]/td/table/tbody/tr[3]/td[1]/table/tbody/tr[2]/td[2]')))
                data_admissao_nome = data_admissao.text
                data_desligamento = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                                    '/html/body/div/table/tbody/tr/td[1]/table/tbody/tr[1]/td[2]/table/tbody/tr[1]/td/table/tbody/tr[3]/td[1]/table/tbody/tr[3]/td[2]')))
                data_desligamento_nome = data_desligamento.text
                browser.switch_to.default_content()
                idetalhe = browser.find_element(By.CSS_SELECTOR, '#iDetalhe')
                browser.switch_to.frame(idetalhe)
                reserva = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                          '/html/body/table/tbody/tr/td/table/tbody/tr[1]/td/table/tbody/tr/td[7]/a/img')))
                reserva.click()
                perfil = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                         '/html/body/table/tbody/tr/td/table/tbody/tr[1]/td/table/tbody/tr/td[8]/table/tbody/tr[1]/td/table/tbody/tr[2]/td[2]/a/img')))
                perfil.click()
                ipcinco = browser.find_element(By.CSS_SELECTOR, '#p5')
                browser.switch_to.frame(ipcinco)
                ipquatrob = browser.find_element(By.CSS_SELECTOR, '#p4b')
                browser.switch_to.frame(ipquatrob)
                tabela = browser.find_elements(By.XPATH, '/html/body/div/div/table/tbody/tr/td[2]/table/tbody/tr/td/a')
                counter = len(tabela) + 1

                for j in range(1, counter):
                    tipo = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                    f'/html/body/div/div/table/tbody/tr/td[2]/table/tbody/tr[{j}]/td/a')))
                    tipo_nome = tipo.text
                    tipo.click()
                    browser.switch_to.default_content()
                    idetalhe = browser.find_element(By.CSS_SELECTOR, '#iDetalhe')
                    browser.switch_to.frame(idetalhe)
                    saldo = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                    '/html/body/table/tbody/tr/td/table/tbody/tr[1]/td/table/tbody/tr/td[8]/table/tbody/tr[1]/td/table/tbody/tr[1]/td[2]/a/img')))
                    saldo.click()
                    ipcinco = browser.find_element(By.CSS_SELECTOR, '#p5')
                    browser.switch_to.frame(ipcinco)
                    ipquatroa1 = browser.find_element(By.CSS_SELECTOR, '#p4a')
                    browser.switch_to.frame(ipquatroa1)
                    tamanho = len(os.listdir(directory))
                    download = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                               '/html/body/table/tbody/tr[2]/td/table/tbody/tr/td[2]/div[2]/table/tbody/tr/td[2]/a/img')))
                    download.click()

                    while True:
                        arq = os.listdir(directory)
                        arquivos = len(list((pal for pal in arq if '.xlsx' or '.xls' in pal)))
                        if tamanho < arquivos:
                            sleep(2)
                            files = os.listdir(directory)
                            files.sort(key=lambda x: os.path.getctime(os.path.join(directory, x)))
                            last_file = files[-1]
                            original = directory + '\\' + last_file
                            new = directory + f'\\{i}_{j}.xlsx'
                            os.rename(original, new)
                            break
                    df = pd.read_excel(directory + f'\\{i}_{j}.xlsx')
                    df['perfil'] = tipo_nome
                    df['plano'] = plano_nome
                    df['data_admissao'] = data_admissao_nome
                    df['data_desligamento'] = data_desligamento_nome
                    df.to_excel(directory + f'\\{i}_{j}.xlsx', index=False)
                    browser.switch_to.default_content()
                    idetalhe = browser.find_element(By.CSS_SELECTOR, '#iDetalhe')
                    browser.switch_to.frame(idetalhe)
                    perfil = WebDriverWait(browser, 10).until(EC.presence_of_element_located((By.XPATH,
                    '/html/body/table/tbody/tr/td/table/tbody/tr[1]/td/table/tbody/tr/td[8]/table/tbody/tr[1]/td/table/tbody/tr[2]/td[2]/a/img')))
                    perfil.click()
                    ipcinco = browser.find_element(By.CSS_SELECTOR, '#p5')
                    browser.switch_to.frame(ipcinco)
                    ipquatrob = browser.find_element(By.CSS_SELECTOR, '#p4b')
                    browser.switch_to.frame(ipquatrob)

                df_total = []
                for x in range(1, counter):
                    df = pd.read_excel(directory + f'\\{i}_{x}.xlsx')
                    df = df[1:]
                    df.rename(columns={'Saldo Participante': 'Saldo Participante Real',
                                       'Unnamed: 3': 'Saldo Participante Qtd. (Cotas)',
                                       'Saldo Patrocinadora': 'Saldo Patrocinadora Real',
                                       'Unnamed: 5': 'Saldo Patrocinadora Qtd. (Cotas)',
                                       'Saldo Total': 'Saldo Total Real',
                                       'Unnamed: 7': 'Saldo Total Qtd. (Cotas)'}, inplace=True)
                    df_total.append(df)
                    os.remove(directory + f'\\{i}_{x}.xlsx')
                df_total = pd.concat(df_total)
                df_total = df_total.sort_values(['Data', 'Valor da Cota'])
                df_total['Data'] = pd.to_datetime(df_total['Data'], format="%d/%m/%Y")
                df_total = df_total.sort_values(['Data', 'Valor da Cota'])
                df_total = df_total.reset_index(drop=True)
                df_total['data_admissao'] = pd.to_datetime(df_total['data_admissao'], format="%d/%m/%Y")
                df_total['data_desligamento'] = pd.to_datetime(df_total['data_desligamento'], format="%d/%m/%Y")

                dias = []
                for x in df_total.index:
                    novo = relativedelta(df_total['data_desligamento'][x], df_total['data_admissao'][x])
                    dias.append(novo.years * 12 + novo.months)
                df_total['meses'] = dias

                perc = []
                for x in df_total.index:
                    if df_total['plano'][x] == 'VISAO MULTI':
                        perc.append(60 if df_total['meses'][x] * 0.25 >= 60 else df_total['meses'][x] * 0.25)
                    else:
                        perc.append(3 if df_total['meses'][x] <= 12 else (
                            6 if df_total['meses'][x] > 12 and df_total['meses'][x] <= 24 else (
                            9 if df_total['meses'][x] > 24 and df_total['meses'][x] <= 36 else (
                            12 if df_total['meses'][x] > 36 and df_total['meses'][x] <= 48 else (
                            60 if df_total['meses'][x] > 48 and df_total['meses'][x] <= 60 else (
                            67.5 if df_total['meses'][x] > 60 and df_total['meses'][x] <= 72 else (
                            75 if df_total['meses'][x] > 72 and df_total['meses'][x] <= 84 else (
                            82.5 if df_total['meses'][x] > 84 and df_total['meses'][x] <= 96 else (
                            90 if df_total['meses'][x] > 96 else df_total['meses'][x])))))))))
                df_total['percentual'] = perc

                cota_nova = []
                for x in df_total.index:
                    cota_nova.append(
                        float(''.join(['.' if i == ',' else i for i in list(str(df_total['Valor da Cota'][x]))])))
                df_total['Valor da Cota'] = cota_nova

                lista = []
                for x in df_total.index:
                    lista.append(df_total['Saldo Patrocinadora Qtd. (Cotas)'][x] * df_total['percentual'][x] / 100)

                lista2 = []
                for x in df_total['Saldo Participante Qtd. (Cotas)'].index:
                    lista2.append(df_total['Saldo Participante Qtd. (Cotas)'][x] + lista[x])

                lista3 = []
                for x in range(0, len(lista2)):
                    if x == 0:
                        lista3.append(round(lista2[x], 4))
                    else:
                        if lista2[x - 1] == 0:
                            lista3.append(round(lista2[x] - lista2[x - 2], 4))
                        elif lista2[x] == 0:
                            lista3.append(0)
                        else:
                            lista3.append(round(lista2[x] - lista2[x - 1], 4))
                df_total['part_cotas'] = lista3

                cotas_mudanca = []
                soma_cotas = []
                df_total = df_total.sort_values(['Data', 'Saldo Participante Real'])
                df_total = df_total.reset_index(drop=True)
                if len(df_total['perfil'].unique()) == 1:
                    resultado = sum(df_total['part_cotas'])
                    penultima = df_total['Valor da Cota'].iloc[-1]
                    segunda = pd.DataFrame([round(resultado, 2),
                            locale.currency(round(resultado * penultima, 2), symbol=True, grouping=True,
                            international=False)], index=['Cotas', 'Real'], columns=['Valor'])

                    tempo = []
                    for y in df_total['Data']:
                        tempo.append(relativedelta(datetime.today(), y).years)
                    df_total['tempo em anos'] = tempo

                    aliquota = []
                    for k in df_total['tempo em anos']:
                        if k >= 10:
                            aliquota.append('10%')
                        elif k >= 8 and k < 10:
                            aliquota.append('15%')
                        elif k >= 6 and k < 8:
                            aliquota.append('20%')
                        elif k >= 4 and k < 6:
                            aliquota.append('25%')
                        elif k >= 2 and k < 4:
                            aliquota.append('30%')
                        else:
                            aliquota.append('35%')
                    df_total['aliquota'] = aliquota

                    come_cota = []
                    come = abs(sum(df_total[df_total['part_cotas'] < 0]['part_cotas']))
                    for e in df_total['part_cotas']:
                        if come != 0:
                            if e >= come:
                                come_cota.append(e - come)
                                come = 0
                            else:
                                come_cota.append(-e)
                                come = come - e
                        else:
                            come_cota.append(e)
                    df_total['come_cota'] = come_cota
                    dez = sum(df_total[(df_total['aliquota'] == '10%') & (df_total['come_cota'] > 0)]['come_cota'])
                    quinze = sum(
                        df_total[(df_total['aliquota'] == '15%') & (df_total['come_cota'] > 0)]['come_cota'])
                    vinte = sum(df_total[(df_total['aliquota'] == '20%') & (df_total['come_cota'] > 0)]['come_cota'])
                    vintecinco = sum(
                        df_total[(df_total['aliquota'] == '25%') & (df_total['come_cota'] > 0)]['come_cota'])
                    trinta = sum(df_total[(df_total['aliquota'] == '30%') & (df_total['come_cota'] > 0)]['come_cota'])
                    trintacinco = sum(
                        df_total[(df_total['aliquota'] == '35%') & (df_total['come_cota'] > 0)]['come_cota'])
                    saldo10.append(dez * penultima)
                    saldo15.append(quinze * penultima)
                    saldo20.append(vinte * penultima)
                    saldo25.append(vintecinco * penultima)
                    saldo30.append(trinta * penultima)
                    saldo35.append(trintacinco * penultima)
                    saldo40.append(dez * penultima + quinze * penultima + vinte * penultima + vintecinco * penultima +
                                   trinta * penultima + trintacinco * penultima)
                else:
                    for x in df_total.index:
                        if x != df_total.index[-1] and x != df_total.index[0]:
                            if df_total['perfil'][x] != df_total['perfil'][x + 1] and df_total['Saldo Participante Real'][x] == 0:
                                cotas_mudanca.append(round(df_total['Valor da Cota'][x], 4))
                                cotas_mudanca.append(round(df_total['Valor da Cota'][x + 1], 4))
                                df_total.loc[x, 'part_cotas'] = sum(soma_cotas)
                                soma_cotas = []
                            elif df_total['perfil'][x] != df_total['perfil'][x - 1] and df_total['Saldo Participante Real'][x - 1] == 0:
                                pass
                            else:
                                cotas_mudanca.append(np.nan)
                                soma_cotas.append(round(df_total['part_cotas'][x], 4))
                        else:
                            cotas_mudanca.append(np.nan)
                            soma_cotas.append(round(df_total['part_cotas'][x], 4))
                    df_total['cotas_mudanca'] = cotas_mudanca
                    var = np.nan
                    for x in df_total.index:
                        if x != df_total.index[-1] and x != df_total.index[0]:
                            if (df_total['Data'][x] == df_total['Data'][x + 1] and
                                    df_total['perfil'][x] == df_total['perfil'][x + 1] and
                                    df_total['Saldo Participante Real'][x] == 0 or
                                    df_total['Data'][x] == df_total['Data'][x - 1] and
                                    df_total['perfil'][x] == df_total['perfil'][x - 1] and
                                    df_total['Saldo Participante Real'][x] == 0):
                                var = x
                    df_total = df_total[df_total.index != var].reset_index(drop=True)
                    tag = 0
                    cot_mudanca = list(df_total[~df_total['cotas_mudanca'].isnull()].reset_index(drop=True)['cotas_mudanca'].values)
                    v1 = 0
                    v2 = 1
                    prim = 0
                    sec = 0
                    inde = 0
                    mass = 0
                    posic = 0
                    value = 0
                    coluna_final = 0
                    resultado = 0
                    nova_tx = []
                    mist = []
                    calculo_tx = []
                    for x in df_total.index:
                        if v2 <= len(cot_mudanca):
                            if x != df_total.index[0] and x != df_total.index[-1]:
                                if df_total['perfil'][x] != df_total['perfil'][x + 1] and df_total['Saldo Participante Real'][x] == 0:
                                    if prim == 0:
                                        value = round(df_total['part_cotas'][x + 1], 2)
                                        mist.append(value)
                                        df_total.loc[x, 'part_cotas'] = sum(df_total['part_cotas'][:x])
                                        df_total.loc[x + 1, 'part_cotas'] = sum(nova_tx)
                                        nova_tx.append(sum(nova_tx))
                                        calculo_tx = nova_tx.copy()
                                        contador = list(range(0, len(calculo_tx)))
                                        posic = 2
                                        for m in contador:
                                            v11 = 2
                                            v22 = 3
                                            for n in range(posic, len(cot_mudanca), 2):
                                                calculo_tx[m] = calculo_tx[m] * cot_mudanca[v11] / cot_mudanca[v22]
                                                v11 = v11 + 2
                                                v22 = v22 + 2
                                            if m == contador[-1]:
                                                calculo_tx[m] = 0
                                                calculo_tx.append(0)
                                        posic += 2
                                        prim = 1
                                        inde = x + 1
                                        v1 = v1 + 2
                                        v2 = v2 + 2
                                    else:
                                        value = round(df_total['part_cotas'][x + 1], 2)
                                        mist.append(value)
                                        df_total.loc[x, 'part_cotas'] = (sum(df_total['part_cotas'][inde:inde + len(
                                            nova_tx[inde:])]) + df_total['part_cotas'][inde - 1] -
                                                                         df_total['part_cotas'][inde] + mist[tag])
                                        df_total.loc[x + 1, 'part_cotas'] = (sum(nova_tx[mass:]) +
                                                                             (df_total['part_cotas'][inde - 1] -
                                                                              df_total['part_cotas'][inde] +
                                                                              mist[tag]) * cot_mudanca[v1] /
                                                                             cot_mudanca[v2])
                                        nova_tx.append(sum(nova_tx[mass:]) + (df_total['part_cotas'][inde - 1] -
                                                                              df_total['part_cotas'][inde] + mist[
                                                                                  tag]) * cot_mudanca[v1] /
                                                       cot_mudanca[v2])
                                        proxima_tx = ([(df_total['part_cotas'][inde - 1] -
                                                        df_total['part_cotas'][inde] + mist[tag]) * cot_mudanca[v1] /
                                                       cot_mudanca[v2]] + nova_tx[mass:])
                                        contador = list(range(0, len(proxima_tx)))
                                        for m in contador:
                                            v11 = v1 + 2
                                            v22 = v2 + 2
                                            for n in range(posic, len(cot_mudanca), 2):
                                                proxima_tx[m] = proxima_tx[m] * cot_mudanca[v11] / cot_mudanca[v22]
                                                v11 = v11 + 2
                                                v22 = v22 + 2
                                            if m == contador[-1]:
                                                proxima_tx[m] = 0
                                                proxima_tx.append(0)
                                        posic += 2
                                        calculo_tx = (calculo_tx + proxima_tx).copy()
                                        tag += 1
                                        inde = x + 1
                                        v1 = v1 + 2
                                        v2 = v2 + 2
                                        mass = x - 1
                                elif df_total['perfil'][x - 1] != df_total['perfil'][x] and df_total['Saldo Participante Real'][x - 1] == 0:
                                    nova_tx.append(value)
                                    mass = x + 1
                                else:
                                    nova_tx.append(df_total['part_cotas'][x] * cot_mudanca[v1] / cot_mudanca[v2])
                            else:
                                nova_tx.append(df_total['part_cotas'][x] * cot_mudanca[v1] / cot_mudanca[v2])
                        else:
                            if sec == 0:
                                df_total.loc[x, 'part_cotas'] = sum(calculo_tx)
                                resultado = sum(calculo_tx + [df_total['part_cotas'][x - 1] -
                                                              df_total['part_cotas'][x] + mist[tag]] + list(
                                    df_total['part_cotas'][x + 1:]))
                                coluna_final = list(
                                    calculo_tx + [
                                        df_total['part_cotas'][x - 1] - df_total['part_cotas'][x] + mist[tag]] +
                                    list(df_total['part_cotas'][x + 1:]))
                                sec = 1

                    final_coluna = []
                    for z in range(0, len(coluna_final)):
                        if coluna_final[z - 1] == 0 and coluna_final[z] == 0:
                            pass
                        else:
                            final_coluna.append(coluna_final[z])
                    penultima = df_total['Valor da Cota'].iloc[-1]
                    segunda = pd.DataFrame([list(df_total['Data']), final_coluna],
                                           index=['Data', 'Cotas Calculadas']).transpose()
                    terceira = pd.DataFrame([round(resultado, 2), locale.currency(round(resultado *
                                penultima, 2), symbol=True, grouping=True, international=False)],
                                index=['Cotas', 'Real'], columns=['Valor'])

                    tempo = []
                    for y in segunda['Data']:
                        tempo.append(relativedelta(datetime.today(), y).years)
                    segunda['tempo em anos'] = tempo

                    aliquota = []
                    for k in segunda['tempo em anos']:
                        if k >= 10:
                            aliquota.append('10%')
                        elif k >= 8 and k < 10:
                            aliquota.append('15%')
                        elif k >= 6 and k < 8:
                            aliquota.append('20%')
                        elif k >= 4 and k < 6:
                            aliquota.append('25%')
                        elif k >= 2 and k < 4:
                            aliquota.append('30%')
                        else:
                            aliquota.append('35%')
                    segunda['aliquota'] = aliquota

                    come_cota = []
                    come = abs(sum(segunda[segunda['Cotas Calculadas'] < 0]['Cotas Calculadas']))
                    for e in segunda['Cotas Calculadas']:
                        if come != 0:
                            if e >= come:
                                come_cota.append(e - come)
                                come = 0
                            else:
                                come_cota.append(-e)
                                come = come - e
                        else:
                            come_cota.append(e)
                    segunda['come_cota'] = come_cota

                    dez = sum(segunda[(segunda['aliquota'] == '10%') & (segunda['come_cota'] > 0)]['come_cota'])
                    quinze = sum(segunda[(segunda['aliquota'] == '15%') & (segunda['come_cota'] > 0)]['come_cota'])
                    vinte = sum(segunda[(segunda['aliquota'] == '20%') & (segunda['come_cota'] > 0)]['come_cota'])
                    vintecinco = sum(segunda[(segunda['aliquota'] == '25%') & (segunda['come_cota'] > 0)]['come_cota'])
                    trinta = sum(segunda[(segunda['aliquota'] == '30%') & (segunda['come_cota'] > 0)]['come_cota'])
                    trintacinco = sum(segunda[(segunda['aliquota'] == '35%') & (segunda['come_cota'] > 0)]['come_cota'])

                    saldo10.append(dez * penultima)
                    saldo15.append(quinze * penultima)
                    saldo20.append(vinte * penultima)
                    saldo25.append(vintecinco * penultima)
                    saldo30.append(trinta * penultima)
                    saldo35.append(trintacinco * penultima)
                    saldo40.append(dez * penultima + quinze * penultima + vinte * penultima + vintecinco * penultima +
                                   trinta * penultima + trintacinco * penultima)

            except Exception:
                saldo10.append(0)
                saldo15.append(0)
                saldo20.append(0)
                saldo25.append(0)
                saldo30.append(0)
                saldo35.append(0)
                saldo40.append(0)

        browser.switch_to.default_content()
        sleep(1)
        browser.quit()
        df = pd.DataFrame([saldo10, saldo15, saldo20, saldo25, saldo30, saldo35, saldo40],
             index=['saldo10', 'saldo15', 'saldo20', 'saldo25', 'saldo30', 'saldo35', 'saldo40']).transpose()
        return df

    def aviso(self):
        # Janela que gera os avisos
        aviso_janela = tk.Toplevel()
        p = PhotoImage(file='Base//logo.png')

        # Janela
        aviso_janela.iconphoto(False, p)
        aviso_janela.title("Resgate")
        aviso_janela.config(width=300, height=200, bg='#ffffff')
        aviso_janela.resizable(width=False, height=False)

        # Botão
        botao_aviso = tk.Button(aviso_janela, text="Fechar", command=aviso_janela.destroy)
        botao_aviso.config(bg='#ffffff', activebackground="#e6e6e6", activeforeground="Black", width=10)
        botao_aviso.place(x=105, y=150)

        # Label
        label_aviso = tk.Label(aviso_janela, text=str(self.primeiro_aviso) + '\n' + str(self.segundo_aviso))
        label_aviso.config(font=("Courier", 10), bg='#ffffff')
        label_aviso.place(x=50, y=60)


Main()
